#!/usr/bin/env python3
"""Structural QA for fallback template PPTX outputs.

This is not a visual renderer. It is a deterministic safety check for the
Gorden-style fallback route when LibreOffice/PowerPoint rendering is unavailable
or unstable. It checks the PPTX can be opened, slide count matches the build
report, warnings are surfaced, and obvious sample-template scientific residue is
not present in visible text.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SAMPLE_RESIDUE_PATTERNS = [
    "NJU120",
    "36R",
    "36\u5143\u73af",
    "\u6cb8\u77f3",
    "\u5b54\u9053",
    "\u7845\u539f\u5b50",
]


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def slide_texts(slide) -> list[str]:
    texts = []
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
    return texts


def audit(pptx: Path, report_path: Path | None = None) -> dict[str, Any]:
    prs = Presentation(pptx)
    result: dict[str, Any] = {
        "$schema": "academic-fallback-structural-audit/v1",
        "pptx": str(pptx),
        "slide_count": len(prs.slides),
        "issues": [],
        "warnings": [],
        "slides": [],
        "status": "passed",
    }

    report = None
    if report_path:
        report = json.loads(report_path.read_text(encoding="utf-8"))
        expected_count = len(report.get("selected_source_slides", []))
        if expected_count and expected_count != len(prs.slides):
            result["issues"].append(f"slide count mismatch: pptx={len(prs.slides)}, report={expected_count}")
        for key in ("overflow_warnings", "ellipsis_warnings", "image_frame_warnings", "untouched_replaceable_slots", "limitations"):
            items = report.get(key) or []
            if items:
                result["warnings"].append({"kind": key, "count": len(items)})

    for index, slide in enumerate(prs.slides, 1):
        texts = slide_texts(slide)
        residue = [
            {"pattern": pattern, "text": text[:120]}
            for text in texts
            for pattern in SAMPLE_RESIDUE_PATTERNS
            if pattern in text
        ]
        if residue:
            result["issues"].append(f"slide {index} contains possible sample scientific residue")
        result["slides"].append(
            {
                "slide_no": index,
                "shape_count": len(list(iter_shapes(slide.shapes))),
                "text_count": len(texts),
                "possible_sample_residue": residue,
            }
        )

    if result["issues"]:
        result["status"] = "failed"
    elif result["warnings"]:
        result["status"] = "needs_review"
    return result


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--build-report", type=Path)
    parser.add_argument("--out", type=Path)
    parser.add_argument("--fail-on-review", action="store_true")
    args = parser.parse_args()

    result = audit(args.pptx, args.build_report)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(args.out)
    else:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    if args.fail_on_review and result["status"] != "passed":
        return 3
    return 0 if result["status"] != "failed" else 2


if __name__ == "__main__":
    raise SystemExit(main())
