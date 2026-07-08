#!/usr/bin/env python3
"""Build a no-Image2 fallback PPTX by editing inherited template slots.

This script is intentionally narrow. It mirrors the useful contract of
template-editing skills such as gorden-ppt-skill: choose source slides, edit
addressed slots, preserve inherited text formatting and image frames, then
prune the template deck to the selected slide order.

It does not draw a new layout. If a requested change cannot be mapped to an
inherited slot, the caller should fix the edit plan or choose another source
slide.
"""

from __future__ import annotations

import argparse
import json
import math
import sys
from pathlib import Path
from typing import Any

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from prepare_fallback_figure import trim_background


DEFAULT_TEMPLATE = Path("academic-slide-minimalist/assets/sample-literature-report.pptx")
DEFAULT_MANIFEST = Path("academic-slide-minimalist/references/sample-template-slot-manifest.json")


def visual_width(text: str) -> float:
    width = 0.0
    for char in text:
        if "\u4e00" <= char <= "\u9fff" or "\u3000" <= char <= "\u303f" or "\uff00" <= char <= "\uffef":
            width += 1.0
        elif char == " ":
            width += 0.35
        elif char.isascii():
            width += 0.5
        else:
            width += 0.8
    return width


def check_overflow(text: str, capacity: dict[str, Any] | None) -> tuple[bool, str]:
    if not capacity or capacity.get("capacity_unknown"):
        return True, ""
    chars_per_line = capacity.get("chars_per_line")
    max_lines = capacity.get("max_lines")
    max_chars = capacity.get("max_chars")
    if not chars_per_line or not max_lines or not max_chars:
        return True, ""
    total = visual_width(text.replace("\n", ""))
    if total <= max_chars:
        return True, ""
    needed = sum(max(1, math.ceil(visual_width(part) / chars_per_line)) for part in (text.split("\n") or [text]))
    message = (
        f"visual_width={total:.0f} > max_chars={max_chars}; "
        f"estimated_lines={needed}, max_lines={max_lines}, font_size_pt={capacity.get('font_size_pt')}"
    )
    if capacity.get("autofit"):
        return True, "AUTOFIT " + message
    return False, message


def has_truncation_ellipsis(text: str) -> bool:
    return text.rstrip().endswith(("...", "..", "\u2026", "\u2026\u2026", "\u7b49\u7b49", "\u7b49\u3002"))


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def find_shape(shapes, shape_id: int):
    for shape in iter_shapes(shapes):
        if shape.shape_id == shape_id:
            return shape
    return None


def load_manifest(path: Path) -> dict[str, Any]:
    manifest = json.loads(path.read_text(encoding="utf-8"))
    pages = {int(page["source_slide"]): page for page in manifest.get("pages", [])}
    slots: dict[tuple[int, str], dict[str, Any]] = {}
    for page in pages.values():
        slide_no = int(page["source_slide"])
        for slot in page.get("text_slots", []) + page.get("image_slots", []) + page.get("object_slots", []):
            slots[(slide_no, slot["slot_id"])] = slot
    manifest["_pages_by_slide"] = pages
    manifest["_slots_by_id"] = slots
    return manifest


def shape_path_for_slot(slot: dict[str, Any] | None) -> list[int]:
    if not slot:
        return []
    return list((slot.get("address") or {}).get("shape_path") or [])


def normalize_plan(plan: dict[str, Any]) -> list[dict[str, Any]]:
    slide_map = plan.get("slide_map")
    if not slide_map:
        selected = plan.get("selected_slides") or []
        slide_map = [
            {
                "output_slide": index + 1,
                "source_slide": source_slide,
                "reuse_mode": "template-prune",
                "edits": [],
            }
            for index, source_slide in enumerate(selected)
        ]
    if not slide_map:
        raise ValueError("fallback edit plan must contain slide_map or selected_slides")
    normalized = []
    for index, item in enumerate(slide_map, 1):
        source_slide = int(item["source_slide"])
        output_slide = int(item.get("output_slide", index))
        normalized.append({**item, "source_slide": source_slide, "output_slide": output_slide})
    return normalized


def replace_text(shape, address: dict[str, Any], new_text: str, expected_text: str | None, strict: bool) -> dict[str, Any]:
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"shape_id {shape.shape_id} has no text frame")
    paragraphs = shape.text_frame.paragraphs
    paragraph_index = int(address.get("paragraph", 0))
    if paragraph_index >= len(paragraphs):
        raise ValueError(f"paragraph {paragraph_index} out of range for shape_id {shape.shape_id}")
    paragraph = paragraphs[paragraph_index]
    runs = list(paragraph.runs)
    if not runs:
        before = paragraph.text
        if expected_text is not None and strict and before != expected_text:
            raise ValueError(f"expected text mismatch: have {before!r}, expected {expected_text!r}")
        paragraph.text = new_text
        return {"before": before, "after": new_text, "mode": "paragraph-text"}

    run_index = address.get("run")
    if run_index is None:
        before = "".join(run.text for run in runs)
        if expected_text is not None and strict and before != expected_text:
            raise ValueError(f"expected text mismatch: have {before!r}, expected {expected_text!r}")
        runs[0].text = new_text
        for run in runs[1:]:
            run.text = ""
        for extra_paragraph in paragraphs[paragraph_index + 1 :]:
            for run in extra_paragraph.runs:
                run.text = ""
        return {"before": before, "after": new_text, "mode": "paragraph-run0"}

    run_index = int(run_index)
    if run_index >= len(runs):
        raise ValueError(f"run {run_index} out of range for shape_id {shape.shape_id}")
    before = runs[run_index].text
    if expected_text is not None and strict and before != expected_text:
        raise ValueError(f"expected text mismatch: have {before!r}, expected {expected_text!r}")
    runs[run_index].text = new_text
    return {"before": before, "after": new_text, "mode": "run"}


def contain_geometry(image_path: Path, left: int, top: int, width: int, height: int) -> tuple[int, int, int, int]:
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    frame_ratio = width / height
    if image_ratio >= frame_ratio:
        new_width = width
        new_height = int(width / image_ratio)
    else:
        new_height = height
        new_width = int(height * image_ratio)
    new_left = int(left + (width - new_width) / 2)
    new_top = int(top + (height - new_height) / 2)
    return new_left, new_top, new_width, new_height


def union_frame(slots: list[dict[str, Any]]) -> dict[str, int]:
    frames = [slot.get("frame", {}).get("emu", {}) for slot in slots]
    frames = [frame for frame in frames if frame.get("width") and frame.get("height")]
    if not frames:
        raise ValueError("cannot compute union frame without image slot geometry")
    left = min(int(frame["left"]) for frame in frames)
    top = min(int(frame["top"]) for frame in frames)
    right = max(int(frame["left"]) + int(frame["width"]) for frame in frames)
    bottom = max(int(frame["top"]) + int(frame["height"]) for frame in frames)
    return {"left": left, "top": top, "width": right - left, "height": bottom - top}


def prepare_image(source: Path, prepared_dir: Path) -> tuple[Path, dict[str, Any]]:
    prepared_dir.mkdir(parents=True, exist_ok=True)
    out = prepared_dir / f"{source.stem}_trimmed.png"
    with Image.open(source) as image:
        original_size = image.size
        prepared, meta = trim_background(image)
    prepared.save(out)
    entry = {
        "source": str(source),
        "prepared": str(out),
        "operation": "trim_background",
        "original_size": list(original_size),
        "prepared_size": list(prepared.size),
        "bbox": meta["bbox"],
        "bbox_removed": meta["removed"],
        "changed": meta["changed"],
        "scientific_content_changed": False,
    }
    return out, entry


def replace_picture_contain(slide, shape, image_path: Path, target_frame: dict[str, int] | None = None) -> dict[str, Any]:
    if target_frame:
        left = int(target_frame["left"])
        top = int(target_frame["top"])
        width = int(target_frame["width"])
        height = int(target_frame["height"])
        frame_source = "adapted"
    else:
        left, top, width, height = int(shape.left), int(shape.top), int(shape.width), int(shape.height)
        frame_source = "inherited"
    new_left, new_top, new_width, new_height = contain_geometry(image_path, left, top, width, height)
    parent = shape._element.getparent()
    index = parent.index(shape._element)
    parent.remove(shape._element)
    new_picture = slide.shapes.add_picture(str(image_path), new_left, new_top, width=new_width, height=new_height)
    new_element = new_picture._element
    new_element.getparent().remove(new_element)
    parent.insert(index, new_element)
    occupancy = (new_width * new_height) / max(1, width * height)
    return {
        "fit": "contain",
        "frame_source": frame_source,
        "original_frame": {"left": left, "top": top, "width": width, "height": height},
        "inserted_frame": {"left": new_left, "top": new_top, "width": new_width, "height": new_height},
        "frame_occupancy": round(occupancy, 4),
    }


def delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def prune_slides(prs, selected: list[int]) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    total = len(slide_ids)
    zero_based = [slide - 1 for slide in selected]
    for slide_index in zero_based:
        if slide_index < 0 or slide_index >= total:
            raise ValueError(f"selected slide {slide_index + 1} out of range 1..{total}")
    new_order = [slide_ids[index] for index in zero_based]
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for slide_id in new_order:
        slide_id_list.append(slide_id)


def apply_plan(args) -> dict[str, Any]:
    plan = json.loads(args.plan.read_text(encoding="utf-8"))
    manifest = load_manifest(args.manifest)
    slide_map = normalize_plan(plan)
    selected_sources = [item["source_slide"] for item in slide_map]
    duplicate_sources = sorted({slide for slide in selected_sources if selected_sources.count(slide) > 1})
    if duplicate_sources:
        raise ValueError(
            "Repeated source slides require true slide cloning, which python-pptx does not expose reliably. "
            f"Choose distinct source slides for now: {duplicate_sources}"
        )

    prs = Presentation(args.template)
    report: dict[str, Any] = {
        "$schema": "academic-fallback-build-report/v1",
        "template": str(args.template),
        "manifest": str(args.manifest),
        "plan": str(args.plan),
        "output": str(args.output),
        "image2_backend_used": False,
        "build_route": "template-prune-and-inherited-slot-edit",
        "selected_source_slides": selected_sources,
        "text_edits": [],
        "image_edits": [],
        "delete_edits": [],
        "overflow_warnings": [],
        "ellipsis_warnings": [],
        "image_frame_warnings": [],
        "untouched_replaceable_slots": [],
        "prepared_figures": [],
        "limitations": [],
    }

    slots_by_id = manifest["_slots_by_id"]
    pages_by_slide = manifest["_pages_by_slide"]
    touched_slots: set[tuple[int, str]] = set()
    touched_shape_paths: dict[int, list[list[int]]] = {}

    for item in slide_map:
        source_slide = item["source_slide"]
        slide = prs.slides[source_slide - 1]
        for edit in item.get("edits", []):
            slot_id = edit.get("slot_id")
            slot = slots_by_id.get((source_slide, slot_id)) if slot_id else None
            address = dict(edit.get("address") or (slot or {}).get("address") or {})
            shape_id = address.get("shape_id")
            if shape_id is None:
                raise ValueError(f"edit on source slide {source_slide} missing slot_id/address: {edit}")
            shape = find_shape(slide.shapes, int(shape_id))
            if shape is None:
                raise ValueError(f"source slide {source_slide}: shape_id {shape_id} not found")

            action = edit.get("action")
            if action == "delete":
                delete_shape(shape)
                if slot_id:
                    touched_slots.add((source_slide, slot_id))
                    touched_shape_paths.setdefault(source_slide, []).append(shape_path_for_slot(slot))
                report["delete_edits"].append({"source_slide": source_slide, "slot_id": slot_id, "shape_id": shape_id})
                continue

            if "new_text" in edit:
                expected = edit.get("expected_text")
                if expected is None and slot:
                    expected = slot.get("current_text")
                result = replace_text(shape, address, str(edit["new_text"]), expected, args.strict)
                if slot_id:
                    touched_slots.add((source_slide, slot_id))
                    touched_shape_paths.setdefault(source_slide, []).append(shape_path_for_slot(slot))
                capacity = (slot or {}).get("capacity")
                fits, message = check_overflow(str(edit["new_text"]), capacity)
                entry = {
                    "source_slide": source_slide,
                    "slot_id": slot_id,
                    "shape_id": shape_id,
                    "result": result,
                }
                report["text_edits"].append(entry)
                if message and not fits:
                    report["overflow_warnings"].append({**entry, "warning": message})
                if has_truncation_ellipsis(str(edit["new_text"])):
                    report["ellipsis_warnings"].append(entry)
                continue

            if "new_image" in edit:
                source_image = Path(edit["new_image"])
                if not source_image.is_absolute():
                    source_image = args.plan.parent / source_image
                prepared_image, prepared_entry = prepare_image(source_image, args.prepared_dir)
                target_frame = None
                if edit.get("frame_scope") == "all_image_slots":
                    target_frame = union_frame(pages_by_slide[source_slide].get("image_slots", []))
                elif edit.get("target_frame"):
                    target_frame = {key: int(value) for key, value in edit["target_frame"].items()}
                picture_result = replace_picture_contain(slide, shape, prepared_image, target_frame)
                if slot_id:
                    touched_slots.add((source_slide, slot_id))
                    touched_shape_paths.setdefault(source_slide, []).append(shape_path_for_slot(slot))
                prepared_entry.update({"source_slide": source_slide, "slot_id": slot_id})
                report["prepared_figures"].append(prepared_entry)
                report["image_edits"].append(
                    {
                        "source_slide": source_slide,
                        "slot_id": slot_id,
                        "shape_id": shape_id,
                        "source_image": str(source_image),
                        "prepared_image": str(prepared_image),
                        "figure_profile": edit.get("figure_profile"),
                        "fit_strategy": edit.get("fit_strategy", "contain"),
                        "frame_scope": edit.get("frame_scope", "slot"),
                        "result": picture_result,
                    }
                )
                if picture_result["frame_occupancy"] < float(edit.get("min_frame_occupancy", 0.45)):
                    report["image_frame_warnings"].append(
                        {
                            "source_slide": source_slide,
                            "slot_id": slot_id,
                            "shape_id": shape_id,
                            "warning": "prepared image occupies a small fraction of the inherited frame",
                            "frame_occupancy": picture_result["frame_occupancy"],
                            "suggestion": "choose a better source slide, use a tighter real crop, or split the figure",
                        }
                    )
                continue

            raise ValueError(f"unsupported edit on source slide {source_slide}: {edit}")

    for source_slide in selected_sources:
        page = pages_by_slide.get(source_slide, {})
        for slot_id in page.get("delete_or_replace_objects", []):
            key = (source_slide, slot_id)
            if key not in touched_slots:
                slot = slots_by_id.get(key)
                shape_id = (slot or {}).get("address", {}).get("shape_id")
                if shape_id is not None and any(shape_id in path for path in touched_shape_paths.get(source_slide, [])):
                    continue
                report["untouched_replaceable_slots"].append({"source_slide": source_slide, "slot_id": slot_id})

    prune_slides(prs, selected_sources)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)

    prepared_manifest = {
        "$schema": "academic-fallback-prepared-figures/v1",
        "figures": report["prepared_figures"],
    }
    args.prepared_manifest.parent.mkdir(parents=True, exist_ok=True)
    args.prepared_manifest.write_text(json.dumps(prepared_manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    report["prepared_figure_manifest"] = str(args.prepared_manifest)

    if duplicate_sources:
        report["limitations"].append("repeated source slides are not supported in this first stable builder")
    warning_count = (
        len(report["overflow_warnings"])
        + len(report["ellipsis_warnings"])
        + len(report["image_frame_warnings"])
        + len(report["untouched_replaceable_slots"])
        + len(report["limitations"])
    )
    report["status"] = "needs_review" if warning_count else "passed"
    report["warning_count"] = warning_count
    return report


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("plan", type=Path, help="fallback_edit_plan.json")
    parser.add_argument("output", type=Path, help="output fallback PPTX")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST)
    parser.add_argument("--report", type=Path)
    parser.add_argument("--prepared-dir", type=Path)
    parser.add_argument("--prepared-manifest", type=Path)
    parser.add_argument("--strict", action="store_true", help="fail on expected_text mismatch")
    parser.add_argument("--fail-on-warnings", action="store_true", help="return non-zero when the build report needs review")
    args = parser.parse_args()

    if args.report is None:
        args.report = args.output.with_suffix(".build-report.json")
    if args.prepared_dir is None:
        args.prepared_dir = args.output.parent / "prepared_figures"
    if args.prepared_manifest is None:
        args.prepared_manifest = args.output.parent / "prepared_figure_manifest.json"

    try:
        report = apply_plan(args)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2

    args.report.parent.mkdir(parents=True, exist_ok=True)
    args.report.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(args.output)
    print(args.report)
    if args.fail_on_warnings and report["status"] != "passed":
        print(f"ERROR: build completed but report status is {report['status']}", file=sys.stderr)
        return 3
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
