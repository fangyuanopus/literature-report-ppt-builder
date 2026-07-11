#!/usr/bin/env python3
"""Verify that an exact-template output preserves its inherited visual system."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


SKILL_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = SKILL_ROOT / "assets" / "sample-literature-report.pptx"
DEFAULT_MANIFEST = SKILL_ROOT / "references" / "sample-template-slot-manifest.json"
ACCENT_RGB = {"8B0D18", "A20F18", "A30D18"}


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def find_shape(shapes, shape_id: int):
    return next((shape for shape in iter_shapes(shapes) if shape.shape_id == shape_id), None)


def geometry(shape) -> tuple[int, int, int, int, int]:
    return (
        int(shape.left),
        int(shape.top),
        int(shape.width),
        int(shape.height),
        int(getattr(shape, "rotation", 0) or 0),
    )


def element_xml(element) -> bytes:
    return etree.tostring(element, with_tail=False)


def load_manifest(path: Path) -> tuple[dict[int, dict[str, Any]], dict[tuple[int, str], dict[str, Any]]]:
    manifest = json.loads(path.read_text(encoding="utf-8"))
    pages = {int(page["source_slide"]): page for page in manifest.get("pages", [])}
    slots: dict[tuple[int, str], dict[str, Any]] = {}
    for slide_no, page in pages.items():
        for slot in page.get("text_slots", []) + page.get("image_slots", []) + page.get("object_slots", []):
            slots[(slide_no, slot["slot_id"])] = slot
    return pages, slots


def normalize_slide_map(plan: dict[str, Any]) -> list[dict[str, Any]]:
    slide_map = list(plan.get("slide_map") or [])
    if not slide_map:
        raise ValueError("exact-template fidelity audit requires slide_map")
    slide_map.sort(key=lambda item: int(item["output_slide"]))
    expected = list(range(1, len(slide_map) + 1))
    actual = [int(item["output_slide"]) for item in slide_map]
    if actual != expected:
        raise ValueError(f"output_slide values must be contiguous and ordered: {actual}")
    return slide_map


def inherited_style_run(paragraph, style_name: str, style_from_run: int | None):
    runs = list(paragraph.runs)
    if not runs:
        return None
    if style_from_run is not None:
        index = int(style_from_run)
        return runs[index] if 0 <= index < len(runs) else None
    if style_name == "base":
        return runs[0]
    if style_name == "emphasis":
        for run in runs:
            color = run.font.color
            try:
                rgb = str(color.rgb) if color.type else None
            except (AttributeError, ValueError):
                rgb = None
            if rgb and rgb.upper() in ACCENT_RGB:
                return run
    return None


def run_properties_xml(run) -> bytes:
    properties = run._r.rPr
    return b"" if properties is None else element_xml(properties)


def compare_master_or_layout(source, target, label: str, issues: list[str]) -> None:
    if source.name != target.name:
        issues.append(f"{label} name changed: {source.name!r} -> {target.name!r}")
    source_tree = source._element.cSld.spTree
    target_tree = target._element.cSld.spTree
    if element_xml(source_tree) != element_xml(target_tree):
        issues.append(f"{label} shape tree changed")
    source_background = source._element.cSld.bg
    target_background = target._element.cSld.bg
    if (source_background is None) != (target_background is None):
        issues.append(f"{label} background presence changed")
    elif source_background is not None and element_xml(source_background) != element_xml(target_background):
        issues.append(f"{label} background changed")


def audit(
    output: Path,
    plan_path: Path,
    template_path: Path,
    manifest_path: Path,
    build_report_path: Path | None,
) -> dict[str, Any]:
    plan = json.loads(plan_path.read_text(encoding="utf-8"))
    slide_map = normalize_slide_map(plan)
    _, slots = load_manifest(manifest_path)
    template = Presentation(template_path)
    result_deck = Presentation(output)
    build_report = (
        json.loads(build_report_path.read_text(encoding="utf-8")) if build_report_path else None
    )
    result: dict[str, Any] = {
        "$schema": "academic-exact-template-fidelity-audit/v1",
        "output": str(output),
        "template": str(template_path),
        "plan": str(plan_path),
        "issues": [],
        "slides": [],
        "status": "passed",
    }
    issues: list[str] = result["issues"]

    if (result_deck.slide_width, result_deck.slide_height) != (template.slide_width, template.slide_height):
        issues.append("slide size differs from template")
    if len(result_deck.slides) != len(slide_map):
        issues.append(f"slide count mismatch: output={len(result_deck.slides)}, plan={len(slide_map)}")
        result["status"] = "failed"
        return result

    source_layouts: dict[str, Any] = {}
    output_layouts: dict[str, Any] = {}
    source_masters: dict[str, Any] = {}
    output_masters: dict[str, Any] = {}

    for item in slide_map:
        output_no = int(item["output_slide"])
        source_no = int(item["source_slide"])
        source_slide = template.slides[source_no - 1]
        target_slide = result_deck.slides[output_no - 1]
        slide_issues: list[str] = []
        source_layout_name = str(source_slide.slide_layout.part.partname)
        target_layout_name = str(target_slide.slide_layout.part.partname)
        source_master_name = str(source_slide.slide_layout.slide_master.part.partname)
        target_master_name = str(target_slide.slide_layout.slide_master.part.partname)
        if target_layout_name != source_layout_name:
            slide_issues.append(
                f"layout changed: source={source_layout_name}, output={target_layout_name}"
            )
        if target_master_name != source_master_name:
            slide_issues.append(
                f"master changed: source={source_master_name}, output={target_master_name}"
            )
        source_layouts[source_layout_name] = source_slide.slide_layout
        output_layouts[target_layout_name] = target_slide.slide_layout
        source_masters[source_master_name] = source_slide.slide_layout.slide_master
        output_masters[target_master_name] = target_slide.slide_layout.slide_master

        touched_top_level_ids: set[int] = set()
        edit_records: list[dict[str, Any]] = []
        top_level_deletes = 0
        for edit in item.get("edits", []):
            slot = slots.get((source_no, edit.get("slot_id")))
            address = dict(edit.get("address") or (slot or {}).get("address") or {})
            path = list(address.get("shape_path") or [])
            shape_id = int(address.get("shape_id") or (path[-1] if path else 0))
            if not shape_id:
                slide_issues.append(f"edit lacks a resolvable shape address: {edit}")
                continue
            top_level_id = int(path[0] if path else shape_id)
            touched_top_level_ids.add(top_level_id)
            if edit.get("action") == "delete" and top_level_id == shape_id:
                top_level_deletes += 1
            edit_records.append(
                {
                    "edit": edit,
                    "shape_id": shape_id,
                    "top_level_id": top_level_id,
                }
            )

        source_top_shapes = {shape.shape_id: shape for shape in source_slide.shapes}
        target_top_shapes = {shape.shape_id: shape for shape in target_slide.shapes}
        expected_top_count = len(source_top_shapes) - top_level_deletes
        if len(target_top_shapes) != expected_top_count:
            slide_issues.append(
                f"unexpected top-level shape count: output={len(target_top_shapes)}, expected={expected_top_count}"
            )

        preserved_count = 0
        for shape_id, source_shape in source_top_shapes.items():
            if shape_id in touched_top_level_ids:
                continue
            target_shape = target_top_shapes.get(shape_id)
            if target_shape is None:
                slide_issues.append(f"preserved template shape {shape_id} is missing")
                continue
            preserved_count += 1
            if geometry(source_shape) != geometry(target_shape):
                slide_issues.append(f"preserved template shape {shape_id} geometry changed")
            if source_shape.shape_type != target_shape.shape_type:
                slide_issues.append(f"preserved template shape {shape_id} type changed")
            if element_xml(source_shape._element) != element_xml(target_shape._element):
                slide_issues.append(f"preserved template shape {shape_id} content/style changed")

        for record in edit_records:
            edit = record["edit"]
            shape_id = record["shape_id"]
            if edit.get("action") == "delete" or "new_image" in edit:
                continue
            source_shape = find_shape(source_slide.shapes, shape_id)
            target_shape = find_shape(target_slide.shapes, shape_id)
            if source_shape is None or target_shape is None:
                slide_issues.append(f"edited text shape {shape_id} is missing")
                continue
            if geometry(source_shape) != geometry(target_shape):
                slide_issues.append(f"edited text shape {shape_id} geometry changed")
            if edit.get("new_paragraphs") is not None:
                specs = edit["new_paragraphs"]
                target_paragraphs = target_shape.text_frame.paragraphs
                source_paragraphs = source_shape.text_frame.paragraphs
                for paragraph_index, spec in enumerate(specs):
                    if paragraph_index >= len(target_paragraphs) or paragraph_index >= len(source_paragraphs):
                        slide_issues.append(f"shape {shape_id} paragraph {paragraph_index} is missing")
                        continue
                    if isinstance(spec, str):
                        expected_runs = [{"text": spec, "style": "base"}]
                    else:
                        expected_runs = spec.get("runs") or []
                    target_runs = list(target_paragraphs[paragraph_index].runs)
                    if len(target_runs) != len(expected_runs):
                        slide_issues.append(
                            f"shape {shape_id} paragraph {paragraph_index} run count changed unexpectedly"
                        )
                        continue
                    for run_index, (target_run, run_spec) in enumerate(zip(target_runs, expected_runs)):
                        if target_run.text != str(run_spec.get("text", "")):
                            slide_issues.append(
                                f"shape {shape_id} paragraph {paragraph_index} run {run_index} text mismatch"
                            )
                        source_run = inherited_style_run(
                            source_paragraphs[paragraph_index],
                            str(run_spec.get("style", "base")),
                            run_spec.get("style_from_run"),
                        )
                        if source_run is None:
                            slide_issues.append(
                                f"shape {shape_id} paragraph {paragraph_index} has no requested inherited style"
                            )
                        elif run_properties_xml(source_run) != run_properties_xml(target_run):
                            slide_issues.append(
                                f"shape {shape_id} paragraph {paragraph_index} run {run_index} style changed"
                            )
            elif "new_text" in edit:
                if target_shape.text_frame.text.rstrip() != str(edit["new_text"]).rstrip():
                    slide_issues.append(f"shape {shape_id} replacement text mismatch")
                source_paragraphs = source_shape.text_frame.paragraphs
                target_paragraphs = target_shape.text_frame.paragraphs
                fallback_source = next((paragraph for paragraph in source_paragraphs if paragraph.runs), None)
                for paragraph_index, expected_text in enumerate(str(edit["new_text"]).split("\n")):
                    if paragraph_index >= len(target_paragraphs):
                        slide_issues.append(f"shape {shape_id} paragraph {paragraph_index} is missing")
                        continue
                    source_paragraph = source_paragraphs[paragraph_index]
                    inherited_run = inherited_style_run(source_paragraph, "base", None)
                    if inherited_run is None and fallback_source is not None:
                        inherited_run = inherited_style_run(fallback_source, "base", None)
                    target_runs = list(target_paragraphs[paragraph_index].runs)
                    if len(target_runs) != 1 or target_runs[0].text != expected_text:
                        slide_issues.append(
                            f"shape {shape_id} paragraph {paragraph_index} plain replacement runs changed unexpectedly"
                        )
                    elif inherited_run is None:
                        slide_issues.append(
                            f"shape {shape_id} paragraph {paragraph_index} has no inherited base style"
                        )
                    elif run_properties_xml(inherited_run) != run_properties_xml(target_runs[0]):
                        slide_issues.append(
                            f"shape {shape_id} paragraph {paragraph_index} plain replacement style changed"
                        )

        if slide_issues:
            issues.extend(f"slide {output_no}: {issue}" for issue in slide_issues)
        result["slides"].append(
            {
                "output_slide": output_no,
                "source_slide": source_no,
                "source_layout": source_layout_name,
                "source_master": source_master_name,
                "preserved_top_level_shapes_checked": preserved_count,
                "issues": slide_issues,
            }
        )

    for part_name, source_layout in source_layouts.items():
        target_layout = output_layouts.get(part_name)
        if target_layout is None:
            issues.append(f"used layout missing from output: {part_name}")
        else:
            compare_master_or_layout(source_layout, target_layout, f"layout {part_name}", issues)
    for part_name, source_master in source_masters.items():
        target_master = output_masters.get(part_name)
        if target_master is None:
            issues.append(f"used master missing from output: {part_name}")
            continue
        compare_master_or_layout(source_master, target_master, f"master {part_name}", issues)
        source_theme = next(
            (rel.target_part for rel in source_master.part.rels.values() if rel.reltype == RT.THEME),
            None,
        )
        target_theme = next(
            (rel.target_part for rel in target_master.part.rels.values() if rel.reltype == RT.THEME),
            None,
        )
        if source_theme is None or target_theme is None or source_theme.blob != target_theme.blob:
            issues.append(f"master theme changed: {part_name}")

    if build_report is not None:
        report_sources = [int(value) for value in build_report.get("selected_source_slides", [])]
        plan_sources = [int(item["source_slide"]) for item in slide_map]
        if report_sources != plan_sources:
            issues.append(f"build report source order differs from plan: {report_sources} vs {plan_sources}")
        for edit in build_report.get("image_edits", []):
            image_result = edit.get("result") or {}
            if image_result.get("frame_source") != "inherited":
                issues.append(
                    f"image {edit.get('slot_id')} did not use its inherited frame: {image_result.get('frame_source')}"
                )
            original = image_result.get("original_frame") or {}
            inserted = image_result.get("inserted_frame") or {}
            if original and inserted:
                within = (
                    inserted["left"] >= original["left"]
                    and inserted["top"] >= original["top"]
                    and inserted["left"] + inserted["width"] <= original["left"] + original["width"]
                    and inserted["top"] + inserted["height"] <= original["top"] + original["height"]
                )
                if not within:
                    issues.append(f"image {edit.get('slot_id')} extends outside its inherited frame")

    if issues:
        result["status"] = "failed"
    result["summary"] = {
        "slide_count": len(result_deck.slides),
        "used_layout_count": len(source_layouts),
        "used_master_count": len(source_masters),
        "issue_count": len(issues),
    }
    return result


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("output", type=Path)
    parser.add_argument("--plan", type=Path, required=True)
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST)
    parser.add_argument("--build-report", type=Path)
    parser.add_argument("--out", type=Path)
    parser.add_argument("--fail-on-review", action="store_true")
    args = parser.parse_args()
    try:
        result = audit(args.output, args.plan, args.template, args.manifest, args.build_report)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2
    payload = json.dumps(result, ensure_ascii=False, indent=2) + "\n"
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(payload, encoding="utf-8")
        print(args.out)
    else:
        print(payload, end="")
    if args.fail_on_review and result["status"] != "passed":
        return 3
    return 0 if result["status"] == "passed" else 2


if __name__ == "__main__":
    raise SystemExit(main())
