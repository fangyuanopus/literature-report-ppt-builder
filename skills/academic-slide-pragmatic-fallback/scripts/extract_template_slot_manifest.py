#!/usr/bin/env python3
"""Extract a reusable slot manifest from the bundled sample PPTX.

The fallback PPTX route should edit inherited template objects rather than
redrawing approximate layouts. This helper inspects the bundled sample deck and
records text/image slots with stable-enough addresses for the current template.
"""

from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


EMU_PER_INCH = 914400
EMU_PER_CM = 360000
PX_PER_INCH = 96


ROLE_BY_SLIDE = {
    1: "cover",
    2: "background",
    3: "background",
    4: "background",
    5: "method",
    6: "method",
    7: "method",
    8: "result",
    9: "result",
    10: "result",
    11: "result",
    12: "result",
    13: "result",
    14: "result",
    15: "application",
    16: "application",
    17: "application",
    18: "summary",
    19: "summary",
    20: "ending",
}

USE_FOR_BY_ROLE = {
    "cover": "paper title, Chinese subtitle, presenter/date metadata",
    "background": "research background, unresolved problem, prior work limitation, or paper positioning",
    "method": "research strategy, workflow, architecture, method overview, or experimental design",
    "result": "main evidence slide with real paper figures and short interpretation",
    "application": "application, performance, generalization, or comparison evidence",
    "summary": "takeaways, limitations, cautious conclusion, or presentation summary",
    "ending": "closing statement only",
}

SELECTION_GUIDANCE_BY_LAYOUT = {
    "cover": "Use once at the beginning. Keep metadata concise.",
    "ending": "Use only as a closing slide. Do not force dense content into it.",
    "text-only logic": "Use for background, transition, or summary when no figure is required.",
    "dominant wide figure": "Use for one large figure or dense evidence that needs width.",
    "figure-left text-right": "Use for one dominant figure with concise interpretation.",
    "text-left figure-right": "Use for one dominant figure with concise interpretation.",
    "two-figure evidence": "Use when two figures need direct comparison or sequential evidence.",
    "multi-figure evidence": "Use only when each panel remains readable; otherwise split slides.",
    "single-figure evidence": "Use for one compact figure with a short caption or callout.",
}


def emu_to_inches(value: int) -> float:
    return round(value / EMU_PER_INCH, 4)


def emu_to_cm(value: int) -> float:
    return round(value / EMU_PER_CM, 2)


def emu_to_px(value: int) -> int:
    return round(value / EMU_PER_INCH * PX_PER_INCH)


def frame(shape) -> dict:
    return {
        "emu": {
            "left": int(shape.left),
            "top": int(shape.top),
            "width": int(shape.width),
            "height": int(shape.height),
        },
        "inches": {
            "left": emu_to_inches(shape.left),
            "top": emu_to_inches(shape.top),
            "width": emu_to_inches(shape.width),
            "height": emu_to_inches(shape.height),
        },
        "cm": {
            "left": emu_to_cm(shape.left),
            "top": emu_to_cm(shape.top),
            "width": emu_to_cm(shape.width),
            "height": emu_to_cm(shape.height),
        },
        "px_1280x720": {
            "left": emu_to_px(shape.left),
            "top": emu_to_px(shape.top),
            "width": emu_to_px(shape.width),
            "height": emu_to_px(shape.height),
        },
    }


def visual_width(text: str) -> float:
    total = 0.0
    for char in text:
        if "\u4e00" <= char <= "\u9fff" or "\u3000" <= char <= "\u303f" or "\uff00" <= char <= "\uffef":
            total += 1.0
        elif char == " ":
            total += 0.35
        elif char.isascii():
            total += 0.5
        else:
            total += 0.8
    return total


def estimate_capacity(shape, font_size_pt: float | None) -> dict:
    width_cm = shape.width / EMU_PER_CM
    height_cm = shape.height / EMU_PER_CM
    size = font_size_pt or 16
    # Mirrors the spirit of gorden-ppt-skill's compute_capacity.py: visual-width
    # capacity is advisory, calibrated with a little slack, never a truncation rule.
    usable_w_pt = max(0.0, width_cm - 0.25) * 28.3465
    usable_h_pt = max(0.0, height_cm - 0.08) * 28.3465
    chars_per_line = max(1, int(usable_w_pt / max(size, 1)))
    max_lines = max(1, int(usable_h_pt / max(size, 1)))
    return {
        "box_cm": [round(width_cm, 2), round(height_cm, 2)],
        "font_size_pt": round(size * 2) / 2,
        "chars_per_line": chars_per_line,
        "max_lines": max_lines,
        "max_chars": int(chars_per_line * max_lines * 1.2),
        "wrap": True,
        "autofit": False,
        "capacity_unknown": False,
    }


def first_font_size(shape) -> float | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                return run.font.size.pt
    return None


def iter_shapes(shapes, path=()):
    for shape in shapes:
        current_path = (*path, shape.shape_id)
        yield shape, current_path
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes, current_path)


def text_role(slide_no: int, shape, text: str) -> str:
    top_in = shape.top / EMU_PER_INCH
    left_in = shape.left / EMU_PER_INCH
    width_in = shape.width / EMU_PER_INCH
    stripped = text.strip()
    if slide_no == 1:
        if top_in < 2.0:
            return "cover title"
        if "汇报" in stripped:
            return "cover metadata"
    if slide_no == 20:
        return "ending statement"
    if top_in < 1.1 and width_in > 8:
        return "slide title"
    if top_in < 0.65 and width_in < 1.8:
        return "navigation label"
    if stripped.isdigit() and top_in > 6.8:
        return "page number"
    if left_in > 10.5 and top_in > 6.0:
        return "figure label"
    if stripped.lower().startswith("fig") or stripped.startswith("图") or stripped.startswith("表"):
        return "figure caption"
    return "sample scientific text"


def is_template_chrome_shape(shape) -> bool:
    top_in = shape.top / EMU_PER_INCH
    left_in = shape.left / EMU_PER_INCH
    width_in = shape.width / EMU_PER_INCH
    height_in = shape.height / EMU_PER_INCH
    if top_in < 0.75:
        return True
    if top_in > 6.75:
        return True
    if width_in > 12.5 and height_in > 6.8:
        return True
    if width_in > 12.0 and height_in < 0.25:
        return True
    if left_in < 0.15 and width_in < 0.3:
        return True
    return False


def layout_type(slide_no: int, text_slots: list[dict], image_slots: list[dict]) -> str:
    if slide_no == 1:
        return "cover"
    if slide_no == 20:
        return "ending"
    if not image_slots:
        return "text-only logic"
    if len(image_slots) >= 3:
        return "multi-figure evidence"
    if len(image_slots) == 2:
        return "two-figure evidence"
    image = image_slots[0]["frame"]["inches"]
    if image["width"] > 7:
        return "dominant wide figure"
    if image["left"] < 3 and image["width"] > 4:
        return "figure-left text-right"
    if image["left"] > 6:
        return "text-left figure-right"
    return "single-figure evidence"


def selection_guidance(layout: str) -> str:
    return SELECTION_GUIDANCE_BY_LAYOUT.get(layout, "Use only when the inherited slots match the slide brief.")


def extract(template: Path) -> dict:
    prs = Presentation(template)
    slide_w_px = emu_to_px(prs.slide_width)
    slide_h_px = emu_to_px(prs.slide_height)
    manifest = {
        "$schema": "academic-fallback-template-slots/v1",
        "template_pptx": "assets/sample-literature-report.pptx",
        "slide_count": len(prs.slides),
        "skip_pages": [],
        "page_roles": {},
        "slide_size": {
            "emu": {"width": int(prs.slide_width), "height": int(prs.slide_height)},
            "inches": {"width": emu_to_inches(prs.slide_width), "height": emu_to_inches(prs.slide_height)},
            "cm": {"width": emu_to_cm(prs.slide_width), "height": emu_to_cm(prs.slide_height)},
            "px": {"width": slide_w_px, "height": slide_h_px},
            "aspect": "16:9",
        },
        "theme": {
            "palette": ["#a20f18", "#1d1d1f", "#5f6268", "#f4f2ef", "#ffffff"],
            "style": "red-black-gray academic literature report",
            "source": "derived from bundled sample deck",
        },
        "rules": [
            "Duplicate mapped source slides before editing.",
            "Replace inherited slots; do not redraw approximate layouts.",
            "Preserve inherited slot position, size, typography, crop, and frame treatment.",
            "Replace sample scientific content only with current paper/SI/user-provided content.",
        ],
        "pages": [],
    }
    page_roles: dict[str, list[int]] = {}

    for slide_no, slide in enumerate(prs.slides, 1):
        text_slots = []
        image_slots = []
        object_slots = []
        keep_objects = []
        delete_or_replace = []

        for shape, shape_path in iter_shapes(slide.shapes):
            shape_type = str(shape.shape_type).split(".")[-1]
            base = {
                "shape_id": shape.shape_id,
                "shape_name": shape.name,
                "shape_type": shape_type,
                "address": {"shape_id": shape.shape_id, "shape_path": list(shape_path)},
                "frame": frame(shape),
            }
            if getattr(shape, "has_text_frame", False):
                text = shape.text_frame.text.strip()
                if text:
                    role = text_role(slide_no, shape, text)
                    font_size = first_font_size(shape)
                    slot = {
                        **base,
                        "slot_id": f"s{slide_no:02d}_text_{shape.shape_id}",
                        "kind": "text",
                        "role": role,
                        "current_text": text,
                        "editable": role not in {"navigation label", "page number"},
                        "preserve_style": True,
                        "capacity": estimate_capacity(shape, font_size),
                        "visual_width": round(visual_width(text), 2),
                    }
                    text_slots.append(slot)
                    if role in {"navigation label", "page number"}:
                        keep_objects.append(slot["slot_id"])
                    else:
                        delete_or_replace.append(slot["slot_id"])
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slot = {
                    **base,
                    "slot_id": f"s{slide_no:02d}_image_{shape.shape_id}",
                    "kind": "image",
                    "role": "sample scientific figure" if slide_no not in {1, 20} else "template image",
                    "editable": True,
                    "preserve_frame": True,
                    "expected_source": "sample deck image; replace with current paper/SI/user figure",
                }
                image_slots.append(slot)
                delete_or_replace.append(slot["slot_id"])
            elif not getattr(shape, "has_text_frame", False) and shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                if slide_no not in {1, 20} and not is_template_chrome_shape(shape):
                    slot = {
                        **base,
                        "slot_id": f"s{slide_no:02d}_object_{shape.shape_id}",
                        "kind": "object",
                        "role": "sample scientific/decorative object; delete unless explicitly kept",
                        "editable": True,
                        "preserve_frame": False,
                        "expected_source": "sample deck object; delete or replace for current paper",
                    }
                    object_slots.append(slot)
                    delete_or_replace.append(slot["slot_id"])

        role = ROLE_BY_SLIDE.get(slide_no, "content")
        page_roles.setdefault(role, []).append(slide_no)
        layout = layout_type(slide_no, text_slots, image_slots)
        page = {
            "source_slide": slide_no,
            "role": role,
            "layout_type": layout,
            "use_for": USE_FOR_BY_ROLE.get(role, "content slide"),
            "selection_guidance": selection_guidance(layout),
            "slot_summary": {
                "editable_text_slots": sum(1 for slot in text_slots if slot.get("editable")),
                "image_slots": len(image_slots),
                "object_slots": len(object_slots),
                "replaceable_objects": len(set(delete_or_replace)),
            },
            "text_slots": text_slots,
            "image_slots": image_slots,
            "object_slots": object_slots,
            "keep_objects": sorted(set(keep_objects)),
            "delete_or_replace_objects": sorted(set(delete_or_replace)),
            "hazards": [],
        }
        if not image_slots and page["role"] in {"result", "application"}:
            page["hazards"].append("text-only source slide; choose only when no real figure is required")
        if any(slot["shape_type"] == "GROUP" for slot in text_slots + image_slots + object_slots):
            page["hazards"].append("contains grouped objects")
        if object_slots:
            page["hazards"].append("contains non-picture vector/scientific objects that must be deleted or explicitly kept")
        manifest["pages"].append(page)
    manifest["page_roles"] = page_roles
    assign_type_scale(manifest)
    return manifest


def assign_type_scale(manifest: dict) -> None:
    counts: dict[float, int] = {}
    for page in manifest["pages"]:
        for slot in page["text_slots"]:
            size = slot.get("capacity", {}).get("font_size_pt")
            if size:
                counts[size] = counts.get(size, 0) + 1
    tiers = sorted(counts.items(), key=lambda item: item[0], reverse=True)
    manifest["type_scale"] = [
        {"level": index + 1, "size_pt": size, "count": count}
        for index, (size, count) in enumerate(tiers)
    ]
    level_by_size = {entry["size_pt"]: entry["level"] for entry in manifest["type_scale"]}
    for page in manifest["pages"]:
        for slot in page["text_slots"]:
            size = slot.get("capacity", {}).get("font_size_pt")
            if size:
                slot["level"] = level_by_size[size]


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--template", type=Path, default=Path("academic-slide-minimalist/assets/sample-literature-report.pptx"))
    parser.add_argument("--out", type=Path, default=Path("academic-slide-minimalist/references/sample-template-slot-manifest.json"))
    args = parser.parse_args()

    manifest = extract(args.template)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    args.out.write_text(json.dumps(manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(args.out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
