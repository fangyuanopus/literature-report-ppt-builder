#!/usr/bin/env python3
"""Structural QA for fallback template PPTX outputs.

This is not a visual renderer. It is a deterministic safety check for the
Gorden-style fallback route when LibreOffice/PowerPoint rendering is unavailable
or unstable. It checks the PPTX can be opened, slide count matches the build
report, warnings are surfaced, and obvious sample-template scientific residue is
not present in visible text.
"""

from __future__ import annotations

import argparse
from collections import Counter
import json
import re
import zipfile
from pathlib import Path
from typing import Any

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT


SAMPLE_RESIDUE_PATTERNS = [
    "A 36-ring zeolite",
    "NJU120",
    "36R",
    "36\u5143\u73af",
    "\u6cb8\u77f3",
    "\u5b54\u9053",
    "\u7845\u539f\u5b50",
]

EXACT_TEMPLATE_ALLOWED_MASTER_TEXT = {
    "基本信息",
    "研究背景",
    "研究思路",
    "研究结果",
    "结果分析",
    "总结启发",
}


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def slide_texts(slide) -> list[str]:
    texts = []
    for shape in iter_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                texts.append(text)
    return texts


GENERIC_MASTER_TEXT = {"单击此处编辑母版标题样式", "单击此处编辑母版文本样式", "二级", "三级", "四级", "五级", "‹#›"}


def master_layout_texts(pptx: Path, allowed_text: set[str] | None = None) -> list[dict[str, Any]]:
    """Return visible text stored in master/layout XML, which python-pptx omits."""
    findings: list[dict[str, Any]] = []
    with zipfile.ZipFile(pptx) as archive:
        for name in archive.namelist():
            if not re.fullmatch(r"ppt/(slideMasters|slideLayouts)/[^/]+\.xml", name):
                continue
            xml = archive.read(name).decode("utf-8", errors="ignore")
            texts = [item.strip() for item in re.findall(r"<a:t>(.*?)</a:t>", xml) if item.strip()]
            suspicious = [
                item
                for item in texts
                if item not in GENERIC_MASTER_TEXT
                and item not in (allowed_text or set())
                and not re.fullmatch(r"\d{4}/\d{1,2}/\d{1,2}", item)
            ]
            if suspicious:
                findings.append({"part": name, "texts": suspicious})
    return findings


def xml_child_text(root, local_name: str) -> str | None:
    for element in root.iter():
        if etree.QName(element).localname == local_name:
            return element.text or ""
    return None


def audit(
    pptx: Path,
    report_path: Path | None = None,
    check_masters: bool = False,
    exact_template: bool = False,
) -> dict[str, Any]:
    prs = Presentation(pptx)
    result: dict[str, Any] = {
        "$schema": "academic-fallback-structural-audit/v1",
        "pptx": str(pptx),
        "slide_count": len(prs.slides),
        "issues": [],
        "warnings": [],
        "slides": [],
        "status": "passed",
    }

    with zipfile.ZipFile(pptx) as archive:
        package_names = archive.namelist()
        package_payloads = {
            name: archive.read(name)
            for name in (
                "docProps/core.xml",
                "docProps/app.xml",
                "ppt/presentation.xml",
            )
            if name in package_names
        }
    duplicate_package_parts = sorted(name for name, count in Counter(package_names).items() if count > 1)
    slide_package_parts = sorted(
        name for name in package_names if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
    )
    notes_package_parts = sorted(
        name for name in package_names if re.fullmatch(r"ppt/notesSlides/notesSlide\d+\.xml", name)
    )
    notes_master_parts = sorted(
        name for name in package_names if re.fullmatch(r"ppt/notesMasters/notesMaster\d+\.xml", name)
    )
    slide_master_parts = sorted(
        name for name in package_names if re.fullmatch(r"ppt/slideMasters/slideMaster\d+\.xml", name)
    )
    slide_layout_parts = sorted(
        name for name in package_names if re.fullmatch(r"ppt/slideLayouts/slideLayout\d+\.xml", name)
    )
    presentation_slide_relationships = [
        rel for rel in prs.part.rels.values() if rel.reltype == RT.SLIDE
    ]
    result["package_integrity"] = {
        "duplicate_part_names": duplicate_package_parts,
        "slide_part_count": len(slide_package_parts),
        "notes_slide_part_count": len(notes_package_parts),
        "notes_master_part_count": len(notes_master_parts),
        "slide_master_part_count": len(slide_master_parts),
        "slide_layout_part_count": len(slide_layout_parts),
        "custom_properties_present": "docProps/custom.xml" in package_names,
        "presentation_slide_relationship_count": len(presentation_slide_relationships),
        "visible_slide_count": len(prs.slides),
    }
    if duplicate_package_parts:
        result["issues"].append(f"package contains {len(duplicate_package_parts)} duplicate ZIP part names")
    if len(slide_package_parts) != len(prs.slides):
        result["issues"].append(
            f"package contains hidden/orphan slide parts: package={len(slide_package_parts)}, visible={len(prs.slides)}"
        )
    if len(presentation_slide_relationships) != len(prs.slides):
        result["issues"].append(
            "presentation contains hidden/orphan slide relationships: "
            f"relationships={len(presentation_slide_relationships)}, visible={len(prs.slides)}"
        )
    if notes_package_parts:
        result["issues"].append(f"package retains {len(notes_package_parts)} inherited notes-slide parts")
    if exact_template and notes_master_parts:
        result["issues"].append(f"package retains {len(notes_master_parts)} inherited notes-master parts")
    if exact_template and "docProps/custom.xml" in package_names:
        result["issues"].append("package retains inherited custom document properties")

    used_master_parts = {
        str(slide.slide_layout.slide_master.part.partname).lstrip("/") for slide in prs.slides
    }
    used_layout_parts = {str(slide.slide_layout.part.partname).lstrip("/") for slide in prs.slides}
    extra_master_parts = sorted(set(slide_master_parts) - used_master_parts)
    extra_layout_parts = sorted(set(slide_layout_parts) - used_layout_parts)
    result["package_integrity"]["unused_slide_master_parts"] = extra_master_parts
    result["package_integrity"]["unused_slide_layout_parts"] = extra_layout_parts
    if exact_template and extra_master_parts:
        result["issues"].append(f"package retains {len(extra_master_parts)} unused slide-master parts")
    if exact_template and extra_layout_parts:
        result["issues"].append(f"package retains {len(extra_layout_parts)} unused slide-layout parts")

    core_values: dict[str, str] = {}
    if "docProps/core.xml" in package_payloads:
        core_root = etree.fromstring(package_payloads["docProps/core.xml"])
        for field in ("title", "subject", "creator", "lastModifiedBy", "created", "modified"):
            value = xml_child_text(core_root, field)
            if value is not None:
                core_values[field] = value
    app_values: dict[str, str] = {}
    if "docProps/app.xml" in package_payloads:
        app_root = etree.fromstring(package_payloads["docProps/app.xml"])
        for field in ("Slides", "Notes", "Paragraphs", "Words", "Company"):
            value = xml_child_text(app_root, field)
            if value is not None:
                app_values[field] = value
    result["document_properties"] = {"core": core_values, "extended": app_values}
    if app_values.get("Slides") != str(len(prs.slides)):
        result["issues"].append(
            f"extended properties Slides={app_values.get('Slides')!r}, expected {len(prs.slides)}"
        )
    if app_values.get("Notes") not in {"0", None}:
        result["issues"].append(f"extended properties Notes={app_values.get('Notes')!r}, expected 0")
    if any(pattern in value for value in core_values.values() for pattern in SAMPLE_RESIDUE_PATTERNS):
        result["issues"].append("core document properties contain sample-paper residue")

    presentation_details = {"visible_slide_ids": [], "section_slide_ids": []}
    if "ppt/presentation.xml" in package_payloads:
        presentation_root = etree.fromstring(package_payloads["ppt/presentation.xml"])
        visible_slide_ids = {
            item.get("id")
            for item in presentation_root.xpath("./*[local-name()='sldIdLst']/*[local-name()='sldId']")
            if item.get("id")
        }
        section_slide_ids = {
            item.get("id")
            for item in presentation_root.xpath(".//*[local-name()='sectionLst']//*[local-name()='sldId']")
            if item.get("id")
        }
        presentation_details = {
            "visible_slide_ids": sorted(visible_slide_ids),
            "section_slide_ids": sorted(section_slide_ids),
        }
        stale_section_ids = sorted(section_slide_ids - visible_slide_ids)
        if stale_section_ids:
            result["issues"].append(f"presentation sections reference removed slide IDs: {stale_section_ids}")
        if exact_template and presentation_root.xpath(".//*[local-name()='notesMasterIdLst']"):
            result["issues"].append("presentation retains an inherited notes-master reference")
    result["presentation_structure"] = presentation_details

    report = None
    if report_path:
        report = json.loads(report_path.read_text(encoding="utf-8"))
        expected_count = len(report.get("selected_source_slides", []))
        if expected_count and expected_count != len(prs.slides):
            result["issues"].append(f"slide count mismatch: pptx={len(prs.slides)}, report={expected_count}")
        for key in (
            "overflow_warnings",
            "ellipsis_warnings",
            "image_frame_warnings",
            "image_readability_warnings",
            "untouched_replaceable_slots",
            "limitations",
        ):
            items = report.get(key) or []
            if items:
                result["warnings"].append({"kind": key, "count": len(items)})
        expected_properties = report.get("document_properties") or {}
        property_pairs = {
            "title": core_values.get("title", ""),
            "subject": core_values.get("subject", ""),
            "creator": core_values.get("creator", ""),
            "last_modified_by": core_values.get("lastModifiedBy", ""),
            "company": app_values.get("Company", ""),
        }
        for key, actual in property_pairs.items():
            if key in expected_properties and str(expected_properties[key]) != actual:
                result["issues"].append(
                    f"document property {key} mismatch: have {actual!r}, expected {expected_properties[key]!r}"
                )

    if check_masters:
        allowed_text = EXACT_TEMPLATE_ALLOWED_MASTER_TEXT if exact_template else None
        findings = master_layout_texts(pptx, allowed_text)
        result["master_layout_text"] = findings
        if findings:
            result["warnings"].append(
                {
                    "kind": "master_layout_content",
                    "message": "Template contains non-generic text in masters/layouts; review before cross-domain reuse.",
                    "count": len(findings),
                }
            )

    for index, slide in enumerate(prs.slides, 1):
        referenced_r_ids = {
            value
            for element in slide._element.iter()
            for value in element.attrib.values()
            if isinstance(value, str) and value.startswith("rId")
        }
        unused_relationships = [
            {
                "relationship_id": rel.rId,
                "relationship_type": rel.reltype,
                "target_part": str(getattr(getattr(rel, "target_part", None), "partname", "")),
            }
            for rel in slide.part.rels.values()
            if rel.reltype != RT.SLIDE_LAYOUT and rel.rId not in referenced_r_ids
        ]
        if unused_relationships:
            result["issues"].append(f"slide {index} retains unused content relationships")
        texts = slide_texts(slide)
        residue = [
            {"pattern": pattern, "text": text[:120]}
            for text in texts
            for pattern in SAMPLE_RESIDUE_PATTERNS
            if pattern in text
        ]
        if residue:
            result["issues"].append(f"slide {index} contains possible sample scientific residue")
        result["slides"].append(
            {
                "slide_no": index,
                "shape_count": len(list(iter_shapes(slide.shapes))),
                "text_count": len(texts),
                "unused_content_relationships": unused_relationships,
                "possible_sample_residue": residue,
            }
        )

    if result["issues"]:
        result["status"] = "failed"
    elif result["warnings"]:
        result["status"] = "needs_review"
    return result


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--build-report", type=Path)
    parser.add_argument("--out", type=Path)
    parser.add_argument("--check-masters", action="store_true", help="Report non-generic text embedded in masters/layouts.")
    parser.add_argument(
        "--exact-template",
        action="store_true",
        help="Allow the bundled navigation labels but enforce exact-template package hygiene.",
    )
    parser.add_argument("--fail-on-review", action="store_true")
    args = parser.parse_args()

    result = audit(args.pptx, args.build_report, args.check_masters, args.exact_template)
    if args.out:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(json.dumps(result, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(args.out)
    else:
        print(json.dumps(result, ensure_ascii=False, indent=2))
    if args.fail_on_review and result["status"] != "passed":
        return 3
    return 0 if result["status"] != "failed" else 2


if __name__ == "__main__":
    raise SystemExit(main())
