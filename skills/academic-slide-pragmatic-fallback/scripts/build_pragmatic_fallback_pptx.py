#!/usr/bin/env python3
"""Build clean, editable academic fallback decks in the bundled sample style.

The deck is built from a blank 16:9 presentation. This intentionally avoids
the sample deck's content-bearing master/layout layers while preserving its
white ground, red-black-gray palette, quiet top navigation, title rule, and
footer/page-number rhythm.
"""

from __future__ import annotations

import argparse
import json
import math
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
RED = RGBColor(139, 13, 24)
BLACK = RGBColor(30, 30, 30)
GRAY = RGBColor(105, 105, 105)
LIGHT = RGBColor(247, 247, 247)
BORDER = RGBColor(190, 190, 190)
WHITE = RGBColor(255, 255, 255)
FONT = "Microsoft YaHei"
FORBIDDEN_AUDIENCE_TEXT = ("pragmatic fallback", "editable fallback", "route c", "image2", "模板路径", "构建器")
PT_PER_CM = 28.3465
EMU_PER_CM = 360000
LINE_HEIGHT = 1.2


def set_name(shape, name: str) -> None:
    shape.name = name


def set_east_asia_font(run) -> None:
    rpr = run._r.get_or_add_rPr()
    rpr.set("ea", FONT)


def visual_width(text: str) -> float:
    """Estimate CJK-aware text width in em units, like a template slot linter."""
    width = 0.0
    for char in text:
        if "\u4e00" <= char <= "\u9fff" or "\u3000" <= char <= "\u303f" or "\uff00" <= char <= "\uffef":
            width += 1.0
        elif char == " ":
            width += 0.35
        elif char.isascii():
            width += 0.5
        else:
            width += 0.8
    return width


def assert_text_fits(name: str, text: str, width, height, size: float, wrap: bool) -> None:
    """Reject copy that cannot fit its fixed Route C template slot.

    Route C uses a code-owned but fixed layout.  Treating every text box as a
    slot prevents one long paragraph from silently breaking the visual rhythm.
    """
    if not str(text).strip():
        return
    width_pt = width / EMU_PER_CM * PT_PER_CM
    height_pt = height / EMU_PER_CM * PT_PER_CM
    chars_per_line = max(1, math.floor(width_pt / size))
    max_lines = max(1, math.floor(height_pt / (size * LINE_HEIGHT))) if wrap else 1
    needed_lines = sum(max(1, math.ceil(visual_width(line) / chars_per_line)) for line in str(text).split("\n"))
    if needed_lines > max_lines:
        raise ValueError(
            f"{name} exceeds its fixed layout slot (needs about {needed_lines} lines; "
            f"capacity is {max_lines}). Rewrite or choose a roomier page type."
        )


def add_rect(slide, name, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    set_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_panel(slide, name, x, y, w, h, color=WHITE, rounded=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    set_name(shape, name)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = BORDER
    shape.line.width = Pt(0.8)
    return shape


def add_text(slide, name, text, x, y, w, h, size, color=BLACK, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    assert_text_fits(name, str(text), w, h, size, wrap)
    box = slide.shapes.add_textbox(x, y, w, h)
    set_name(box, name)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = wrap
    frame.margin_left = frame.margin_right = Pt(0)
    frame.margin_top = frame.margin_bottom = Pt(0)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = str(text)
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    set_east_asia_font(run)
    return box


def title_copy(title: str) -> tuple[str, int, float]:
    title = title.strip()
    if len(title) > 54:
        raise ValueError("title exceeds 54 characters; rewrite it before building the slide")
    if len(title) <= 32:
        return title, 30, 0.54
    if "\n" not in title:
        midpoint = len(title) // 2
        split = max(title.rfind("，", 0, midpoint), title.rfind("：", 0, midpoint), title.rfind(" ", 0, midpoint))
        if split < 12:
            split = midpoint
        title = title[: split + 1].rstrip() + "\n" + title[split + 1 :].lstrip()
    return title, 24, 0.84


def validate_copy(item: dict) -> None:
    for value in [item.get("title", ""), item.get("takeaway", "")] + list(item.get("bullets") or []):
        lowered = str(value).lower()
        if any(term in lowered for term in FORBIDDEN_AUDIENCE_TEXT):
            raise ValueError(f"slide {item.get('slide_no')}: internal production text is not audience-facing: {value}")
    bullets = item.get("bullets") or []
    if len(bullets) > 5 or any(len(str(line)) > 58 for line in bullets):
        raise ValueError(f"slide {item.get('slide_no')}: shorten to at most five bullets of 58 characters")


def add_chrome(slide, navigation: list[str], section: str, slide_no: int, footer: str):
    if section not in navigation:
        raise ValueError(f"slide {slide_no}: section must be one of deck.navigation")
    nav_x, nav_y, nav_w, nav_h = Inches(0.16), Inches(0.03), Inches(13.01), Inches(0.38)
    slot_w = nav_w / len(navigation)
    for index, label in enumerate(navigation):
        x = nav_x + slot_w * index
        active = label == section
        add_panel(slide, f"AGENT_NAV_SLOT_{index}", x, nav_y, slot_w, nav_h, RED if active else WHITE)
        add_text(slide, f"AGENT_NAV_LABEL_{index}", label, x, Inches(0.065), slot_w, Inches(0.25), 12, WHITE if active else BLACK, active, PP_ALIGN.CENTER, False)
    if footer:
        add_text(slide, "AGENT_FOOTER", footer, Inches(0.35), Inches(7.12), Inches(9.4), Inches(0.14), 8, GRAY, False, PP_ALIGN.LEFT, False)
    add_text(slide, "AGENT_PAGE_NUMBER", str(slide_no), Inches(12.6), Inches(7.11), Inches(0.28), Inches(0.14), 8, GRAY, False, PP_ALIGN.RIGHT, False)


def add_title(slide, title: str):
    copy, size, height = title_copy(title)
    add_text(slide, "AGENT_TITLE", copy, Inches(0.35), Inches(0.59), Inches(12.35), Inches(height), size, RED, False, PP_ALIGN.LEFT, True)


def add_bullets(slide, bullets: list[str], x, y, w, h, size=18):
    if not bullets:
        return
    text = "\n".join(f"{index}. {line}" for index, line in enumerate(bullets, start=1))
    add_text(slide, "AGENT_BODY", text, x, y, w, h, size, BLACK, False, PP_ALIGN.LEFT, True)


def add_figure(slide, figure: dict, x, y, w, h, slide_no: int, index: int):
    path = Path(figure.get("path", ""))
    if not path.is_file():
        raise FileNotFoundError(f"slide {slide_no}: missing prepared figure: {path}")
    required = ("figure_id", "source_page", "source_label")
    if figure.get("source_type") != "figure_crop" or figure.get("crop_verified") is not True or any(not figure.get(key) for key in required):
        raise ValueError(f"slide {slide_no}: figure needs figure_crop, crop_verified=true, figure_id, source_page, and source_label")
    with Image.open(path) as image:
        ratio = image.width / image.height
    frame_ratio = w / h
    if ratio >= frame_ratio:
        actual_w, actual_h = w, int(w / ratio)
        actual_x, actual_y = x, y + int((h - actual_h) / 2)
    else:
        actual_w, actual_h = int(h * ratio), h
        actual_x, actual_y = x + int((w - actual_w) / 2), y
    picture = slide.shapes.add_picture(str(path), actual_x, actual_y, actual_w, actual_h)
    set_name(picture, f"AGENT_FIGURE_{figure['figure_id']}")
    caption = figure.get("caption") or f"来源：{figure['source_label']}（论文第 {figure['source_page']} 页）"
    add_text(slide, f"AGENT_FIGURE_CAPTION_{index}", caption, x, y + h + Inches(0.02), w, Inches(0.18), 9, GRAY, False, PP_ALIGN.CENTER, False)


def layout_cover(slide, item):
    english_title = item.get("english_title")
    if english_title:
        add_text(slide, "AGENT_COVER_ENGLISH_TITLE", english_title, Inches(0.9), Inches(1.62), Inches(11.55), Inches(0.46), 20, BLACK, False, PP_ALIGN.CENTER, True)
    add_text(slide, "AGENT_COVER_TITLE", item["title"], Inches(0.9), Inches(2.04), Inches(11.55), Inches(0.72), 30, RED, False, PP_ALIGN.CENTER, True)
    add_rect(slide, "AGENT_COVER_DIVIDER", Inches(1.0), Inches(2.82), Inches(11.33), Inches(0.018), RED)
    meta = item.get("metadata") or item.get("bullets") or []
    if meta:
        if len(meta) >= 2:
            add_text(slide, "AGENT_COVER_META_LEFT", str(meta[0]), Inches(1.5), Inches(3.55), Inches(4.4), Inches(0.4), 16, BLACK, False, PP_ALIGN.CENTER, True)
            add_text(slide, "AGENT_COVER_META_RIGHT", str(meta[1]), Inches(7.4), Inches(3.55), Inches(4.4), Inches(0.4), 16, BLACK, False, PP_ALIGN.CENTER, True)
        else:
            add_text(slide, "AGENT_COVER_META", str(meta[0]), Inches(2.0), Inches(3.55), Inches(9.3), Inches(0.4), 16, BLACK, False, PP_ALIGN.CENTER, True)
    if item.get("takeaway"):
        add_text(slide, "AGENT_TAKEAWAY", item["takeaway"], Inches(1.1), Inches(5.4), Inches(11.1), Inches(0.32), 15, RED, False, PP_ALIGN.CENTER, True)
    add_rect(slide, "AGENT_COVER_BOTTOM_BAND", 0, Inches(7.0), SLIDE_W, Inches(0.5), RED)


def layout_text(slide, item):
    add_title(slide, item["title"])
    add_panel(slide, "AGENT_TEXT_PANEL", Inches(0.75), Inches(1.55), Inches(11.85), Inches(4.85), WHITE, True)
    add_bullets(slide, item.get("bullets") or [], Inches(1.05), Inches(1.82), Inches(11.25), Inches(4.3), 20)


def layout_figure_right(slide, item):
    add_title(slide, item["title"])
    figures = item.get("figures") or []
    if len(figures) != 1:
        raise ValueError("figure_right requires exactly one figure")
    add_panel(slide, "AGENT_TEXT_PANEL", Inches(0.45), Inches(1.45), Inches(4.5), Inches(4.9), LIGHT, True)
    add_bullets(slide, item.get("bullets") or [], Inches(0.72), Inches(1.72), Inches(3.95), Inches(4.25), 18)
    add_panel(slide, "AGENT_EVIDENCE_PANEL", Inches(5.15), Inches(1.35), Inches(7.72), Inches(5.08), WHITE, True)
    add_figure(slide, figures[0], Inches(5.35), Inches(1.55), Inches(7.32), Inches(4.5), item["slide_no"], 0)


def layout_figure_wide(slide, item):
    add_title(slide, item["title"])
    figures = item.get("figures") or []
    if len(figures) != 1:
        raise ValueError("figure_wide requires exactly one figure")
    with Image.open(figures[0]["path"]) as image:
        ratio = image.width / image.height
    if ratio < 1.75:
        add_panel(slide, "AGENT_EVIDENCE_PANEL", Inches(0.42), Inches(1.35), Inches(7.72), Inches(5.08), WHITE, True)
        add_figure(slide, figures[0], Inches(0.62), Inches(1.55), Inches(7.32), Inches(4.5), item["slide_no"], 0)
        add_panel(slide, "AGENT_TEXT_PANEL", Inches(8.38), Inches(1.45), Inches(4.5), Inches(4.9), LIGHT, True)
        add_bullets(slide, item.get("bullets") or [], Inches(8.65), Inches(1.72), Inches(3.95), Inches(4.25), 18)
    else:
        add_panel(slide, "AGENT_EVIDENCE_PANEL", Inches(0.42), Inches(1.35), Inches(12.46), Inches(4.42), WHITE, True)
        add_figure(slide, figures[0], Inches(0.62), Inches(1.52), Inches(12.06), Inches(3.88), item["slide_no"], 0)
        add_bullets(slide, item.get("bullets") or [], Inches(0.75), Inches(5.86), Inches(11.8), Inches(0.58), 17)


def layout_process(slide, item):
    add_title(slide, item["title"])
    steps = item.get("steps") or []
    if not 2 <= len(steps) <= 5:
        raise ValueError("process requires 2–5 short steps")
    width = Inches(10.9 / len(steps))
    for index, step in enumerate(steps):
        x = Inches(1.15) + index * width
        add_rect(slide, f"AGENT_STEP_BOX_{index}", x, Inches(3.0), width - Inches(0.16), Inches(1.0), LIGHT)
        add_text(slide, f"AGENT_STEP_TEXT_{index}", step, x + Inches(0.08), Inches(3.15), width - Inches(0.32), Inches(0.7), 18, BLACK, True, PP_ALIGN.CENTER, True)
        if index < len(steps) - 1:
            add_text(slide, f"AGENT_STEP_ARROW_{index}", "→", x + width - Inches(0.08), Inches(3.21), Inches(0.18), Inches(0.4), 20, RED, True, PP_ALIGN.CENTER, False)
    add_bullets(slide, item.get("bullets") or [], Inches(1.1), Inches(4.65), Inches(10.9), Inches(1.15), 18)


def layout_comparison(slide, item):
    add_title(slide, item["title"])
    columns = item.get("columns") or []
    if len(columns) != 2:
        raise ValueError("comparison requires exactly two columns")
    for index, column in enumerate(columns):
        x = Inches(0.95 + index * 5.85)
        add_text(slide, f"AGENT_COLUMN_HEAD_{index}", column.get("heading", ""), x, Inches(1.65), Inches(5.15), Inches(0.35), 21, RED, True, PP_ALIGN.LEFT, True)
        add_bullets(slide, column.get("bullets") or [], x, Inches(2.15), Inches(5.15), Inches(3.7), 18)


def layout_summary(slide, item):
    add_title(slide, item["title"])
    rows = item.get("summary_rows") or []
    if not 2 <= len(rows) <= 5:
        raise ValueError("summary requires 2–5 summary_rows")
    row_h = Inches(4.65 / len(rows))
    for index, row in enumerate(rows):
        y = Inches(1.55) + index * row_h
        add_rect(slide, f"AGENT_SUMMARY_LINE_{index}", Inches(0.95), y, Inches(11.35), Inches(0.012), RED)
        add_text(slide, f"AGENT_SUMMARY_LABEL_{index}", row.get("label", ""), Inches(1.05), y + Inches(0.12), Inches(2.0), row_h - Inches(0.18), 19, RED, True, PP_ALIGN.LEFT, True)
        add_text(slide, f"AGENT_SUMMARY_VALUE_{index}", row.get("value", ""), Inches(3.2), y + Inches(0.12), Inches(8.65), row_h - Inches(0.18), 18, BLACK, False, PP_ALIGN.LEFT, True)


def layout_closing(slide, item):
    add_text(slide, "AGENT_CLOSING", item.get("closing_text", "汇报完毕，敬请批评指正！"), Inches(0.95), Inches(2.95), Inches(11.43), Inches(0.92), 54, BLACK, False, PP_ALIGN.CENTER, True)
    if item.get("metadata"):
        add_text(slide, "AGENT_CLOSING_META", "\n".join(item["metadata"][:3]), Inches(2.0), Inches(4.35), Inches(9.3), Inches(0.55), 13, GRAY, False, PP_ALIGN.CENTER, True)
    add_rect(slide, "AGENT_CLOSING_BOTTOM_BAND", 0, Inches(7.0), SLIDE_W, Inches(0.5), RED)


def infer_layout(item: dict) -> str:
    if item.get("layout_type"):
        return item["layout_type"]
    if item.get("closing_text"):
        return "closing"
    if item.get("slide_no") == 1:
        return "cover"
    if item.get("steps"):
        return "process"
    if item.get("columns"):
        return "comparison"
    if item.get("summary_rows"):
        return "summary"
    figures = item.get("figures") or []
    if len(figures) == 1:
        with Image.open(figures[0]["path"]) as image:
            return "figure_wide" if image.width / image.height >= 1.15 else "figure_right"
    return "text"


LAYOUTS = {"cover": layout_cover, "text": layout_text, "figure_right": layout_figure_right, "figure_wide": layout_figure_wide, "process": layout_process, "comparison": layout_comparison, "summary": layout_summary, "closing": layout_closing}


def build(plan: dict, output: Path) -> dict:
    deck, slides = plan.get("deck") or {}, plan.get("slides") or []
    navigation = deck.get("navigation") or []
    if not (3 <= len(navigation) <= 7) or not slides:
        raise ValueError("deck needs 3–7 navigation labels and at least one slide")
    if len(set(navigation)) != len(navigation) or any(not str(label).strip() for label in navigation):
        raise ValueError("deck.navigation labels must be non-empty and unique")

    slide_numbers = []
    for item in slides:
        try:
            slide_numbers.append(int(item.get("slide_no", 0)))
        except (TypeError, ValueError) as exc:
            raise ValueError("every slide needs a positive integer slide_no") from exc
    expected_numbers = list(range(1, len(slides) + 1))
    if slide_numbers != expected_numbers:
        raise ValueError(f"slide_no values must be sequential 1..{len(slides)} in deck order")

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    footer = deck.get("footer", "")
    lowered_footer = str(footer).lower()
    if any(term in lowered_footer for term in FORBIDDEN_AUDIENCE_TEXT):
        raise ValueError("deck.footer must be bibliographic and cannot expose internal production terms")
    used_layouts = []
    for item in slides:
        item = dict(item)
        item["slide_no"] = int(item["slide_no"])
        if not item.get("title"):
            raise ValueError("every slide needs a title")
        validate_copy(item)
        layout = infer_layout(item)
        if layout not in LAYOUTS:
            raise ValueError(f"slide {item['slide_no']}: unsupported layout_type {layout}")
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        if layout not in ("cover", "closing"):
            add_chrome(slide, navigation, item.get("section") or navigation[0], item["slide_no"], footer)
        LAYOUTS[layout](slide, item)
        if item.get("takeaway") and layout not in ("cover", "closing"):
            add_text(slide, "AGENT_TAKEAWAY", item["takeaway"], Inches(0.95), Inches(6.55), Inches(11.15), Inches(0.3), 16, RED, True, PP_ALIGN.LEFT, True)
        used_layouts.append(layout)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return {"status": "passed", "slide_count": len(slides), "output": str(output), "route": "pragmatic_fallback", "layouts": used_layouts}


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("plan", type=Path)
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument("--report", type=Path)
    args = parser.parse_args()
    report = build(json.loads(args.plan.read_text(encoding="utf-8")), args.out)
    data = json.dumps(report, ensure_ascii=False, indent=2) + "\n"
    if args.report:
        args.report.write_text(data, encoding="utf-8")
    print(data, end="")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
