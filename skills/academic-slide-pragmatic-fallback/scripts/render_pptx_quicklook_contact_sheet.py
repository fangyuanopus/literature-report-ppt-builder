#!/usr/bin/env python3
"""Render a PPTX into per-slide QuickLook thumbnails and a contact sheet.

This is a macOS visual QA helper for fallback template PPTX decks. It uses
QuickLook because LibreOffice rendering is not always available or stable in
the local Codex desktop environment.
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


def split_pptx(pptx: Path, out_dir: Path) -> list[Path]:
    prs = Presentation(pptx)
    slide_count = len(prs.slides)
    singles_dir = out_dir / "single_slide_pptx"
    singles_dir.mkdir(parents=True, exist_ok=True)
    paths = []
    for index in range(slide_count):
        single = Presentation(pptx)
        keep = index
        slide_id_list = single.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
        for delete_index in range(slide_count - 1, -1, -1):
            if delete_index == keep:
                continue
            r_id = slide_id_list[delete_index].rId
            single.part.drop_rel(r_id)
            del slide_id_list[delete_index]
        path = singles_dir / f"slide_{index + 1:02d}.pptx"
        single.save(path)
        paths.append(path)
    return paths


def render_with_quicklook(single_pptx: Path, render_dir: Path) -> Path:
    render_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        ["qlmanage", "-t", "-s", "1600", "-o", str(render_dir), str(single_pptx)],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    candidates = sorted(render_dir.glob(single_pptx.name + "*.png"))
    if not candidates:
        candidates = sorted(render_dir.glob("*.png"), key=lambda p: p.stat().st_mtime, reverse=True)
    if not candidates:
        raise FileNotFoundError(f"QuickLook did not produce a PNG for {single_pptx}")
    return candidates[-1]


def make_contact_sheet(images: list[Path], out: Path, columns: int = 4) -> None:
    thumbs = []
    for index, image_path in enumerate(images, 1):
        image = Image.open(image_path).convert("RGB")
        image.thumbnail((520, 292), Image.LANCZOS)
        tile = Image.new("RGB", (560, 352), "white")
        x = (560 - image.width) // 2
        tile.paste(image, (x, 34))
        draw = ImageDraw.Draw(tile)
        try:
            font = ImageFont.truetype("Arial.ttf", 18)
        except OSError:
            font = ImageFont.load_default()
        draw.text((14, 8), f"Slide {index}", fill=(40, 40, 40), font=font)
        thumbs.append(tile)
    rows = (len(thumbs) + columns - 1) // columns
    sheet = Image.new("RGB", (columns * 560, rows * 352), (238, 238, 238))
    for index, thumb in enumerate(thumbs):
        x = (index % columns) * 560
        y = (index // columns) * 352
        sheet.paste(thumb, (x, y))
    out.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out)


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--out-dir", type=Path, required=True)
    parser.add_argument("--columns", type=int, default=4)
    parser.add_argument("--clean", action="store_true")
    args = parser.parse_args()

    if args.clean and args.out_dir.exists():
        shutil.rmtree(args.out_dir)
    args.out_dir.mkdir(parents=True, exist_ok=True)
    render_dir = args.out_dir / "rendered_slides"
    singles = split_pptx(args.pptx, args.out_dir)
    rendered = [render_with_quicklook(path, render_dir) for path in singles]
    contact_sheet = args.out_dir / "contact_sheet.png"
    make_contact_sheet(rendered, contact_sheet, columns=args.columns)
    print(contact_sheet)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
