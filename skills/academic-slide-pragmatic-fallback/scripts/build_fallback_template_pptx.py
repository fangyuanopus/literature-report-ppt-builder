#!/usr/bin/env python3
"""Build a no-Image2 fallback PPTX by editing inherited template slots.

This script is intentionally narrow. It mirrors the useful contract of
template-editing skills such as gorden-ppt-skill: choose source slides, edit
addressed slots, preserve inherited text formatting and image frames, then
prune the template deck to the selected slide order.

It does not draw a new layout. If a requested change cannot be mapped to an
inherited slot, the caller should fix the edit plan or choose another source
slide.
"""

from __future__ import annotations

import argparse
from collections import Counter
from copy import deepcopy
import json
import math
import os
import re
import sys
import zipfile
from pathlib import Path
from typing import Any

from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn

from prepare_fallback_figure import trim_background


SKILL_ROOT = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = SKILL_ROOT / "assets" / "sample-literature-report.pptx"
DEFAULT_MANIFEST = SKILL_ROOT / "references" / "sample-template-slot-manifest.json"


def visual_width(text: str) -> float:
    width = 0.0
    for char in text:
        if "\u4e00" <= char <= "\u9fff" or "\u3000" <= char <= "\u303f" or "\uff00" <= char <= "\uffef":
            width += 1.0
        elif char == " ":
            width += 0.35
        elif char.isascii():
            # Mixed Latin/CJK captions render wider than the old half-em
            # approximation in the bundled Chinese template.
            width += 0.6
        else:
            width += 0.8
    return width


def check_overflow(text: str, capacity: dict[str, Any] | None) -> tuple[bool, str]:
    if not capacity or capacity.get("capacity_unknown"):
        return True, ""
    chars_per_line = capacity.get("chars_per_line")
    max_lines = capacity.get("max_lines")
    max_chars = capacity.get("max_chars")
    if not chars_per_line or not max_lines or not max_chars:
        return True, ""
    total = visual_width(text.replace("\n", ""))
    if total <= max_chars:
        return True, ""
    needed = sum(max(1, math.ceil(visual_width(part) / chars_per_line)) for part in (text.split("\n") or [text]))
    message = (
        f"visual_width={total:.0f} > max_chars={max_chars}; "
        f"estimated_lines={needed}, max_lines={max_lines}, font_size_pt={capacity.get('font_size_pt')}"
    )
    if capacity.get("autofit"):
        return True, "AUTOFIT " + message
    return False, message


def has_truncation_ellipsis(text: str) -> bool:
    return text.rstrip().endswith(("...", "..", "\u2026", "\u2026\u2026", "\u7b49\u7b49", "\u7b49\u3002"))


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def find_shape(shapes, shape_id: int):
    for shape in iter_shapes(shapes):
        if shape.shape_id == shape_id:
            return shape
    return None


def load_manifest(path: Path) -> dict[str, Any]:
    manifest = json.loads(path.read_text(encoding="utf-8"))
    pages = {int(page["source_slide"]): page for page in manifest.get("pages", [])}
    slots: dict[tuple[int, str], dict[str, Any]] = {}
    for page in pages.values():
        slide_no = int(page["source_slide"])
        for slot in page.get("text_slots", []) + page.get("image_slots", []) + page.get("object_slots", []):
            slots[(slide_no, slot["slot_id"])] = slot
    manifest["_pages_by_slide"] = pages
    manifest["_slots_by_id"] = slots
    return manifest


def shape_path_for_slot(slot: dict[str, Any] | None) -> list[int]:
    if not slot:
        return []
    return list((slot.get("address") or {}).get("shape_path") or [])


def normalize_plan(plan: dict[str, Any]) -> list[dict[str, Any]]:
    slide_map = plan.get("slide_map")
    if not slide_map:
        selected = plan.get("selected_slides") or []
        slide_map = [
            {
                "output_slide": index + 1,
                "source_slide": source_slide,
                "reuse_mode": "template-prune",
                "edits": [],
            }
            for index, source_slide in enumerate(selected)
        ]
    if not slide_map:
        raise ValueError("fallback edit plan must contain slide_map or selected_slides")
    normalized = []
    for index, item in enumerate(slide_map, 1):
        source_slide = int(item["source_slide"])
        output_slide = int(item.get("output_slide", index))
        normalized.append({**item, "source_slide": source_slide, "output_slide": output_slide})
    output_numbers = [item["output_slide"] for item in normalized]
    if len(set(output_numbers)) != len(output_numbers):
        raise ValueError(f"output_slide values must be unique: {output_numbers}")
    expected_numbers = list(range(1, len(normalized) + 1))
    if sorted(output_numbers) != expected_numbers:
        raise ValueError(f"output_slide values must be contiguous 1..{len(normalized)}: {output_numbers}")
    return sorted(normalized, key=lambda item: item["output_slide"])


def comparable_text(text: str | None) -> str | None:
    return text.rstrip() if text is not None else None


def rebuild_paragraph_runs(
    paragraph,
    run_specs: list[dict[str, Any]] | str,
    fallback_run_elements: list[Any] | None = None,
) -> None:
    if isinstance(run_specs, str):
        run_specs = [{"text": run_specs, "style": "base"}]
    source_runs = list(paragraph.runs)
    source_run_elements = [deepcopy(run._r) for run in source_runs]
    if not source_run_elements and fallback_run_elements:
        source_run_elements = [deepcopy(run_element) for run_element in fallback_run_elements]

    base_index = 0
    emphasis_index = None
    for index, run_element in enumerate(source_run_elements):
        colors = run_element.xpath(".//*[local-name()='srgbClr']/@val")
        if any(str(color).upper() in {"8B0D18", "A20F18", "A30D18"} for color in colors):
            emphasis_index = index
            break

    paragraph_element = paragraph._p
    for child in list(paragraph_element):
        if etree.QName(child).localname in {"r", "br", "fld"}:
            paragraph_element.remove(child)

    for run_spec in run_specs:
        text = str(run_spec.get("text", ""))
        style_from_run = run_spec.get("style_from_run")
        if style_from_run is None:
            style_name = run_spec.get("style", "base")
            if style_name == "emphasis":
                if emphasis_index is None:
                    raise ValueError("paragraph has no inherited emphasis run to reuse")
                style_from_run = emphasis_index
            elif style_name == "base":
                style_from_run = base_index
            else:
                raise ValueError(f"unsupported inherited run style: {style_name}")
        style_from_run = int(style_from_run)
        if source_run_elements:
            if style_from_run < 0 or style_from_run >= len(source_run_elements):
                raise ValueError(
                    f"style_from_run {style_from_run} out of range for paragraph with {len(source_run_elements)} runs"
                )
            run_element = deepcopy(source_run_elements[style_from_run])
            run_element.text = text
            end_properties = next(
                (
                    child
                    for child in paragraph_element
                    if etree.QName(child).localname == "endParaRPr"
                ),
                None,
            )
            if end_properties is None:
                paragraph_element.append(run_element)
            else:
                paragraph_element.insert(paragraph_element.index(end_properties), run_element)
        else:
            paragraph.add_run().text = text


def replace_whole_text_box(shape, new_text: str | None, new_paragraphs: list[Any] | None) -> dict[str, Any]:
    paragraphs = shape.text_frame.paragraphs
    fallback_run_elements = next(
        ([deepcopy(run._r) for run in paragraph.runs] for paragraph in paragraphs if paragraph.runs),
        [],
    )
    if new_paragraphs is None:
        paragraph_specs: list[Any] = str(new_text or "").split("\n")
    else:
        paragraph_specs = list(new_paragraphs)
    if len(paragraph_specs) > len(paragraphs):
        raise ValueError(
            f"new content requires {len(paragraph_specs)} paragraphs but inherited text box has {len(paragraphs)}"
        )
    for index, paragraph in enumerate(paragraphs):
        if index >= len(paragraph_specs):
            rebuild_paragraph_runs(paragraph, "", fallback_run_elements)
            continue
        spec = paragraph_specs[index]
        if isinstance(spec, str):
            rebuild_paragraph_runs(paragraph, spec, fallback_run_elements)
        elif isinstance(spec, dict):
            runs = spec.get("runs")
            if not isinstance(runs, list) or not runs:
                raise ValueError("new_paragraphs entries must be strings or objects with a non-empty runs list")
            rebuild_paragraph_runs(paragraph, runs, fallback_run_elements)
        else:
            raise ValueError("new_paragraphs entries must be strings or objects")
    return {
        "after": "\n".join(paragraph.text for paragraph in paragraphs).rstrip("\n"),
        "mode": "whole-text-box-paragraph-aware",
        "paragraph_count": len(paragraph_specs),
    }


def paragraph_specs_text(new_paragraphs: list[Any]) -> str:
    lines = []
    for spec in new_paragraphs:
        if isinstance(spec, str):
            lines.append(spec)
        elif isinstance(spec, dict):
            lines.append("".join(str(run.get("text", "")) for run in spec.get("runs", [])))
        else:
            lines.append(str(spec))
    return "\n".join(lines)


def replace_text(
    shape,
    address: dict[str, Any],
    new_text: str | None,
    expected_text: str | None,
    strict: bool,
    new_paragraphs: list[Any] | None = None,
) -> dict[str, Any]:
    if not getattr(shape, "has_text_frame", False):
        raise ValueError(f"shape_id {shape.shape_id} has no text frame")
    paragraphs = shape.text_frame.paragraphs
    whole_text_box = "paragraph" not in address and "run" not in address
    if whole_text_box:
        before = shape.text_frame.text
        if expected_text is not None and strict and comparable_text(before) != comparable_text(expected_text):
            raise ValueError(f"expected text mismatch: have {before!r}, expected {expected_text!r}")
        result = replace_whole_text_box(shape, new_text, new_paragraphs)
        return {"before": before, **result}
    paragraph_index = int(address.get("paragraph", 0))
    if paragraph_index >= len(paragraphs):
        raise ValueError(f"paragraph {paragraph_index} out of range for shape_id {shape.shape_id}")
    paragraph = paragraphs[paragraph_index]
    runs = list(paragraph.runs)
    if not runs:
        before = paragraph.text
        if expected_text is not None and strict and comparable_text(before) != comparable_text(expected_text):
            raise ValueError(f"expected text mismatch: have {before!r}, expected {expected_text!r}")
        paragraph.text = str(new_text or "")
        return {"before": before, "after": new_text, "mode": "paragraph-text"}

    run_index = address.get("run")
    if run_index is None:
        before = "".join(run.text for run in runs)
        if expected_text is not None and strict and comparable_text(before) != comparable_text(expected_text):
            raise ValueError(f"expected text mismatch: have {before!r}, expected {expected_text!r}")
        runs[0].text = str(new_text or "")
        for run in runs[1:]:
            run.text = ""
        for extra_paragraph in paragraphs[paragraph_index + 1 :]:
            for run in extra_paragraph.runs:
                run.text = ""
        return {"before": before, "after": new_text, "mode": "paragraph-run0"}

    run_index = int(run_index)
    if run_index >= len(runs):
        raise ValueError(f"run {run_index} out of range for shape_id {shape.shape_id}")
    before = runs[run_index].text
    if expected_text is not None and strict and comparable_text(before) != comparable_text(expected_text):
        raise ValueError(f"expected text mismatch: have {before!r}, expected {expected_text!r}")
    runs[run_index].text = str(new_text or "")
    return {"before": before, "after": str(new_text or ""), "mode": "run"}


def contain_geometry(image_path: Path, left: int, top: int, width: int, height: int) -> tuple[int, int, int, int]:
    with Image.open(image_path) as image:
        image_width, image_height = image.size
    image_ratio = image_width / image_height
    frame_ratio = width / height
    if image_ratio >= frame_ratio:
        new_width = width
        new_height = int(width / image_ratio)
    else:
        new_height = height
        new_width = int(height * image_ratio)
    new_left = int(left + (width - new_width) / 2)
    new_top = int(top + (height - new_height) / 2)
    return new_left, new_top, new_width, new_height


def union_frame(slots: list[dict[str, Any]]) -> dict[str, int]:
    frames = [slot.get("frame", {}).get("emu", {}) for slot in slots]
    frames = [frame for frame in frames if frame.get("width") and frame.get("height")]
    if not frames:
        raise ValueError("cannot compute union frame without image slot geometry")
    left = min(int(frame["left"]) for frame in frames)
    top = min(int(frame["top"]) for frame in frames)
    right = max(int(frame["left"]) + int(frame["width"]) for frame in frames)
    bottom = max(int(frame["top"]) + int(frame["height"]) for frame in frames)
    return {"left": left, "top": top, "width": right - left, "height": bottom - top}


def prepare_image(source: Path, prepared_dir: Path) -> tuple[Path, dict[str, Any]]:
    prepared_dir.mkdir(parents=True, exist_ok=True)
    out = prepared_dir / f"{source.stem}_trimmed.png"
    with Image.open(source) as image:
        original_size = image.size
        prepared, meta = trim_background(image)
    prepared.save(out)
    entry = {
        "source": str(source),
        "prepared": str(out),
        "operation": "trim_background",
        "original_size": list(original_size),
        "prepared_size": list(prepared.size),
        "bbox": meta["bbox"],
        "bbox_removed": meta["removed"],
        "changed": meta["changed"],
        "scientific_content_changed": False,
    }
    return out, entry


def replace_picture_contain(slide, shape, image_path: Path, target_frame: dict[str, int] | None = None) -> dict[str, Any]:
    if target_frame:
        left = int(target_frame["left"])
        top = int(target_frame["top"])
        width = int(target_frame["width"])
        height = int(target_frame["height"])
        frame_source = "adapted"
    else:
        left, top, width, height = int(shape.left), int(shape.top), int(shape.width), int(shape.height)
        frame_source = "inherited"
    new_left, new_top, new_width, new_height = contain_geometry(image_path, left, top, width, height)
    parent = shape._element.getparent()
    index = parent.index(shape._element)
    parent.remove(shape._element)
    new_picture = slide.shapes.add_picture(str(image_path), new_left, new_top, width=new_width, height=new_height)
    new_element = new_picture._element
    new_element.getparent().remove(new_element)
    parent.insert(index, new_element)
    occupancy = (new_width * new_height) / max(1, width * height)
    return {
        "fit": "contain",
        "frame_source": frame_source,
        "original_frame": {"left": left, "top": top, "width": width, "height": height},
        "inserted_frame": {"left": new_left, "top": new_top, "width": new_width, "height": new_height},
        "frame_occupancy": round(occupancy, 4),
    }


def image_display_size_px(picture_result: dict[str, Any], slide_width: int, slide_height: int) -> tuple[int, int]:
    inserted = picture_result["inserted_frame"]
    width_px = round(int(inserted["width"]) / max(1, slide_width) * 1280)
    height_px = round(int(inserted["height"]) / max(1, slide_height) * 720)
    return width_px, height_px


def delete_shape(shape) -> None:
    shape._element.getparent().remove(shape._element)


def clean_slide_relationships(slide) -> list[dict[str, Any]]:
    """Remove prior-template payload no longer referenced by slide XML.

    Removing a picture shape does not make python-pptx drop its image
    relationship. Likewise, inherited notes remain reachable even though they
    are not represented by a shape. Both cases can preserve source-paper data
    inside an otherwise clean-looking output package.
    """
    referenced_r_ids = {
        value
        for element in slide._element.iter()
        for value in element.attrib.values()
        if isinstance(value, str) and value.startswith("rId")
    }
    removed: list[dict[str, Any]] = []
    for rel in list(slide.part.rels.values()):
        if rel.reltype == RT.SLIDE_LAYOUT:
            continue
        remove_reason = None
        if rel.reltype == RT.NOTES_SLIDE:
            remove_reason = "inherited-notes"
        elif rel.rId not in referenced_r_ids:
            remove_reason = "unreferenced"
        if remove_reason is None:
            continue
        removed.append(
            {
                "slide_part": str(slide.part.partname),
                "relationship_id": rel.rId,
                "relationship_type": rel.reltype,
                "target_part": str(getattr(getattr(rel, "target_part", None), "partname", "")),
                "reason": remove_reason,
            }
        )
        slide.part.drop_rel(rel.rId)
    return removed


def remove_relationship_list(parent_part, list_element) -> list[dict[str, str]]:
    removed: list[dict[str, str]] = []
    for item in list(list_element):
        relationship_id = item.get(qn("r:id"))
        if not relationship_id:
            continue
        rel = parent_part.rels[relationship_id]
        removed.append(
            {
                "relationship_id": relationship_id,
                "relationship_type": rel.reltype,
                "target_part": str(getattr(getattr(rel, "target_part", None), "partname", "")),
            }
        )
        parent_part.drop_rel(relationship_id)
        list_element.remove(item)
    return removed


def clean_presentation_structure(prs) -> dict[str, Any]:
    """Keep only masters/layouts used by visible slides and remove stale UI state."""
    used_layout_parts = {slide.slide_layout.part for slide in prs.slides}
    used_master_parts = {slide.slide_layout.slide_master.part for slide in prs.slides}
    report: dict[str, Any] = {
        "removed_slide_masters": [],
        "removed_slide_layouts": [],
        "removed_notes_masters": [],
        "removed_sections": 0,
    }

    for master in list(prs.slide_masters):
        if master.part not in used_master_parts:
            continue
        for child in list(master.part._element):
            if etree.QName(child).localname != "sldLayoutIdLst":
                continue
            for item in list(child):
                relationship_id = item.get(qn("r:id"))
                if not relationship_id:
                    continue
                rel = master.part.rels[relationship_id]
                if rel.target_part in used_layout_parts:
                    continue
                report["removed_slide_layouts"].append(str(rel.target_part.partname))
                master.part.drop_rel(relationship_id)
                child.remove(item)

    presentation_element = prs.part._element
    for child in list(presentation_element):
        local_name = etree.QName(child).localname
        if local_name == "sldMasterIdLst":
            for item in list(child):
                relationship_id = item.get(qn("r:id"))
                if not relationship_id:
                    continue
                rel = prs.part.rels[relationship_id]
                if rel.target_part in used_master_parts:
                    continue
                report["removed_slide_masters"].append(str(rel.target_part.partname))
                prs.part.drop_rel(relationship_id)
                child.remove(item)
        elif local_name == "notesMasterIdLst":
            report["removed_notes_masters"].extend(remove_relationship_list(prs.part, child))
            presentation_element.remove(child)

    for section_list in list(presentation_element.xpath(".//*[local-name()='sectionLst']")):
        parent = section_list.getparent()
        parent.remove(section_list)
        report["removed_sections"] += 1
    for extension in list(presentation_element.xpath(".//*[local-name()='ext']")):
        if len(extension) == 0 and not (extension.text or "").strip():
            extension.getparent().remove(extension)
    for extension_list in list(presentation_element.xpath(".//*[local-name()='extLst']")):
        if len(extension_list) == 0:
            extension_list.getparent().remove(extension_list)
    return report


def derive_document_title(plan: dict[str, Any], slide_map: list[dict[str, Any]], manifest: dict[str, Any]) -> str:
    slots_by_id = manifest["_slots_by_id"]
    for item in slide_map:
        if item.get("role") != "cover":
            continue
        source_slide = int(item["source_slide"])
        for edit in item.get("edits", []):
            slot = slots_by_id.get((source_slide, edit.get("slot_id")))
            if slot and slot.get("role") == "cover title" and edit.get("new_text"):
                return str(edit["new_text"])
    for item in slide_map:
        for edit in item.get("edits", []):
            if edit.get("new_text"):
                return str(edit["new_text"])
    return "Academic literature report"


def sanitize_core_properties(prs, plan: dict[str, Any], slide_map: list[dict[str, Any]], manifest: dict[str, Any]) -> dict[str, Any]:
    requested = dict(plan.get("document_properties") or {})
    title = str(requested.get("title") or derive_document_title(plan, slide_map, manifest))
    subject = str(requested.get("subject") or "学术文献汇报")
    author = str(requested.get("creator") or "")
    last_modified_by = str(requested.get("last_modified_by") or "")
    core = prs.core_properties
    core.title = title
    core.subject = subject
    core.author = author
    core.last_modified_by = last_modified_by
    core.keywords = str(requested.get("keywords") or "")
    core.comments = str(requested.get("comments") or "")
    core.category = str(requested.get("category") or "")
    core.language = str(requested.get("language") or "")
    core.content_status = ""
    core.identifier = ""
    core.version = ""
    core.revision = 1
    return {
        "title": title,
        "subject": subject,
        "creator": author,
        "last_modified_by": last_modified_by,
        "company": str(requested.get("company") or ""),
    }


def presentation_text_stats(prs) -> tuple[int, int]:
    paragraph_count = 0
    word_count = 0
    for slide in prs.slides:
        for shape in iter_shapes(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            paragraph_count += len(shape.text_frame.paragraphs)
            text = shape.text_frame.text
            word_count += len(re.findall(r"[A-Za-z0-9]+|[\u3400-\u9fff]", text))
    return paragraph_count, word_count


def sanitize_package_properties(
    pptx: Path,
    *,
    slide_count: int,
    paragraph_count: int,
    word_count: int,
    company: str,
) -> None:
    """Rewrite stale extended/custom properties after python-pptx saves."""
    app_namespace = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
    with zipfile.ZipFile(pptx, "r") as source:
        entries = [(item, source.read(item.filename)) for item in source.infolist()]

    rewritten: list[tuple[zipfile.ZipInfo, bytes]] = []
    for item, payload in entries:
        name = item.filename
        if name == "docProps/custom.xml":
            continue
        if name == "docProps/core.xml":
            root = etree.fromstring(payload)
            removable_core_fields = {
                "created",
                "modified",
                "lastPrinted",
            }
            for element in list(root):
                if etree.QName(element).localname in removable_core_fields:
                    root.remove(element)
            payload = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
        elif name == "docProps/app.xml":
            root = etree.fromstring(payload)
            values = {
                "TotalTime": 0,
                "Words": word_count,
                "Paragraphs": paragraph_count,
                "Slides": slide_count,
                "Notes": 0,
                "HiddenSlides": 0,
            }
            for local_name, value in values.items():
                elements = root.findall(f"{{{app_namespace}}}{local_name}")
                if elements:
                    elements[0].text = str(value)
            for local_name in ("HeadingPairs", "TitlesOfParts", "Manager"):
                for element in root.findall(f"{{{app_namespace}}}{local_name}"):
                    root.remove(element)
            company_elements = root.findall(f"{{{app_namespace}}}Company")
            if company_elements:
                company_elements[0].text = company
            payload = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
        elif name == "_rels/.rels":
            root = etree.fromstring(payload)
            for relationship in list(root):
                if relationship.get("Type", "").endswith("/custom-properties"):
                    root.remove(relationship)
            payload = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
        elif name == "[Content_Types].xml":
            root = etree.fromstring(payload)
            for override in list(root):
                if override.get("PartName") == "/docProps/custom.xml":
                    root.remove(override)
            payload = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
        rewritten.append((item, payload))

    temp_path = pptx.with_suffix(pptx.suffix + ".tmp")
    with zipfile.ZipFile(temp_path, "w") as target:
        for item, payload in rewritten:
            target.writestr(item, payload)
    os.replace(temp_path, pptx)
    with zipfile.ZipFile(pptx, "r") as archive:
        duplicates = [name for name, count in Counter(archive.namelist()).items() if count > 1]
    if duplicates:
        raise ValueError(f"metadata rewrite produced duplicate package parts: {duplicates}")


def prune_slides(prs, selected: list[int]) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)
    total = len(slide_ids)
    zero_based = [slide - 1 for slide in selected]
    for slide_index in zero_based:
        if slide_index < 0 or slide_index >= total:
            raise ValueError(f"selected slide {slide_index + 1} out of range 1..{total}")
    new_order = [slide_ids[index] for index in zero_based]
    retained_relationship_ids = {slide_id.rId for slide_id in new_order}
    for slide_id in slide_ids:
        if slide_id.rId not in retained_relationship_ids:
            # Removing only the <p:sldId> leaves the old slide relationship in
            # presentation.xml.rels. The slide then remains embedded in the
            # package even though PowerPoint does not show it. Drop the
            # relationship as well so prior-paper slides and their media are
            # unreachable and omitted when python-pptx saves the package.
            prs.part.drop_rel(slide_id.rId)
    for slide_id in list(slide_id_list):
        slide_id_list.remove(slide_id)
    for slide_id in new_order:
        slide_id_list.append(slide_id)


def apply_plan(args) -> dict[str, Any]:
    plan = json.loads(args.plan.read_text(encoding="utf-8"))
    manifest = load_manifest(args.manifest)
    slide_map = normalize_plan(plan)
    selected_sources = [item["source_slide"] for item in slide_map]
    duplicate_sources = sorted({slide for slide in selected_sources if selected_sources.count(slide) > 1})
    if duplicate_sources:
        raise ValueError(
            "Repeated source slides require true slide cloning, which python-pptx does not expose reliably. "
            f"Choose distinct source slides for now: {duplicate_sources}"
        )

    prs = Presentation(args.template)
    report: dict[str, Any] = {
        "$schema": "academic-fallback-build-report/v1",
        "template": str(args.template),
        "manifest": str(args.manifest),
        "plan": str(args.plan),
        "output": str(args.output),
        "image2_backend_used": False,
        "build_route": "template-prune-and-inherited-slot-edit",
        "selected_source_slides": selected_sources,
        "text_edits": [],
        "image_edits": [],
        "delete_edits": [],
        "overflow_warnings": [],
        "ellipsis_warnings": [],
        "image_frame_warnings": [],
        "image_readability_warnings": [],
        "untouched_replaceable_slots": [],
        "prepared_figures": [],
        "removed_inherited_relationships": [],
        "presentation_cleanup": {},
        "document_properties": {},
        "limitations": [],
    }

    slots_by_id = manifest["_slots_by_id"]
    pages_by_slide = manifest["_pages_by_slide"]
    touched_slots: set[tuple[int, str]] = set()
    touched_shape_paths: dict[int, list[list[int]]] = {}

    for item in slide_map:
        source_slide = item["source_slide"]
        slide = prs.slides[source_slide - 1]
        for edit in item.get("edits", []):
            slot_id = edit.get("slot_id")
            slot = slots_by_id.get((source_slide, slot_id)) if slot_id else None
            address = dict(edit.get("address") or (slot or {}).get("address") or {})
            shape_id = address.get("shape_id")
            if shape_id is None:
                raise ValueError(f"edit on source slide {source_slide} missing slot_id/address: {edit}")
            shape = find_shape(slide.shapes, int(shape_id))
            if shape is None:
                raise ValueError(f"source slide {source_slide}: shape_id {shape_id} not found")

            action = edit.get("action")
            if action == "delete":
                delete_shape(shape)
                if slot_id:
                    touched_slots.add((source_slide, slot_id))
                    touched_shape_paths.setdefault(source_slide, []).append(shape_path_for_slot(slot))
                report["delete_edits"].append({"source_slide": source_slide, "slot_id": slot_id, "shape_id": shape_id})
                continue

            if "new_text" in edit or "new_paragraphs" in edit:
                expected = edit.get("expected_text")
                if expected is None and slot:
                    expected = slot.get("current_text")
                new_paragraphs = edit.get("new_paragraphs")
                display_text = (
                    str(edit.get("new_text", ""))
                    if new_paragraphs is None
                    else paragraph_specs_text(new_paragraphs)
                )
                result = replace_text(
                    shape,
                    address,
                    edit.get("new_text"),
                    expected,
                    args.strict,
                    new_paragraphs,
                )
                if slot_id:
                    touched_slots.add((source_slide, slot_id))
                    touched_shape_paths.setdefault(source_slide, []).append(shape_path_for_slot(slot))
                capacity = (slot or {}).get("capacity")
                fits, message = check_overflow(display_text, capacity)
                entry = {
                    "source_slide": source_slide,
                    "slot_id": slot_id,
                    "shape_id": shape_id,
                    "result": result,
                }
                report["text_edits"].append(entry)
                if message and not fits:
                    report["overflow_warnings"].append({**entry, "warning": message})
                if has_truncation_ellipsis(display_text):
                    report["ellipsis_warnings"].append(entry)
                continue

            if "new_image" in edit:
                source_image = Path(edit["new_image"])
                if not source_image.is_absolute():
                    source_image = args.plan.parent / source_image
                prepared_image, prepared_entry = prepare_image(source_image, args.prepared_dir)
                fit_strategy = str(edit.get("fit_strategy", "contain"))
                if fit_strategy not in {"contain", "adaptive_contain"}:
                    raise ValueError(f"unsupported fit_strategy {fit_strategy!r} for slot {slot_id}")
                target_frame = None
                if edit.get("frame_scope") == "all_image_slots":
                    target_frame = union_frame(pages_by_slide[source_slide].get("image_slots", []))
                elif edit.get("target_frame"):
                    target_frame = {key: int(value) for key, value in edit["target_frame"].items()}
                if fit_strategy == "adaptive_contain" and target_frame is None:
                    raise ValueError("adaptive_contain requires frame_scope=all_image_slots or target_frame")
                if fit_strategy == "contain" and target_frame is not None:
                    raise ValueError("contain must use the inherited slot frame; use adaptive_contain for a larger frame")
                picture_result = replace_picture_contain(slide, shape, prepared_image, target_frame)
                if slot_id:
                    touched_slots.add((source_slide, slot_id))
                    touched_shape_paths.setdefault(source_slide, []).append(shape_path_for_slot(slot))
                prepared_entry.update({"source_slide": source_slide, "slot_id": slot_id})
                report["prepared_figures"].append(prepared_entry)
                report["image_edits"].append(
                    {
                        "source_slide": source_slide,
                        "slot_id": slot_id,
                        "shape_id": shape_id,
                        "source_image": str(source_image),
                        "prepared_image": str(prepared_image),
                        "figure_profile": edit.get("figure_profile"),
                        "fit_strategy": fit_strategy,
                        "frame_scope": edit.get("frame_scope", "slot"),
                        "result": picture_result,
                    }
                )
                if picture_result["frame_occupancy"] < float(edit.get("min_frame_occupancy", 0.45)):
                    report["image_frame_warnings"].append(
                        {
                            "source_slide": source_slide,
                            "slot_id": slot_id,
                            "shape_id": shape_id,
                            "warning": "prepared image occupies a small fraction of the inherited frame",
                            "frame_occupancy": picture_result["frame_occupancy"],
                            "suggestion": "choose a better source slide, use a tighter real crop, or split the figure",
                        }
                    )
                profile = edit.get("figure_profile") or {}
                if profile.get("dense") or profile.get("kind") == "table":
                    width_px, height_px = image_display_size_px(
                        picture_result,
                        int(prs.slide_width),
                        int(prs.slide_height),
                    )
                    min_width_px = int(edit.get("min_display_width_px", 560))
                    min_height_px = int(edit.get("min_display_height_px", 170))
                    if width_px < min_width_px or height_px < min_height_px:
                        report["image_readability_warnings"].append(
                            {
                                "source_slide": source_slide,
                                "slot_id": slot_id,
                                "shape_id": shape_id,
                                "warning": "dense table is too small for reliable reading at 1280x720",
                                "display_size_px": [width_px, height_px],
                                "minimum_size_px": [min_width_px, min_height_px],
                                "suggestion": "use a wider inherited frame, split the table, or replace it with a more legible real figure",
                            }
                        )
                continue

            raise ValueError(f"unsupported edit on source slide {source_slide}: {edit}")

    for source_slide in selected_sources:
        page = pages_by_slide.get(source_slide, {})
        for slot_id in page.get("delete_or_replace_objects", []):
            key = (source_slide, slot_id)
            if key not in touched_slots:
                slot = slots_by_id.get(key)
                shape_id = (slot or {}).get("address", {}).get("shape_id")
                if shape_id is not None and any(shape_id in path for path in touched_shape_paths.get(source_slide, [])):
                    continue
                report["untouched_replaceable_slots"].append({"source_slide": source_slide, "slot_id": slot_id})

    for source_slide in selected_sources:
        report["removed_inherited_relationships"].extend(clean_slide_relationships(prs.slides[source_slide - 1]))

    prune_slides(prs, selected_sources)
    report["presentation_cleanup"] = clean_presentation_structure(prs)
    report["document_properties"] = sanitize_core_properties(prs, plan, slide_map, manifest)
    paragraph_count, word_count = presentation_text_stats(prs)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    sanitize_package_properties(
        args.output,
        slide_count=len(prs.slides),
        paragraph_count=paragraph_count,
        word_count=word_count,
        company=report["document_properties"]["company"],
    )

    prepared_manifest = {
        "$schema": "academic-fallback-prepared-figures/v1",
        "figures": report["prepared_figures"],
    }
    args.prepared_manifest.parent.mkdir(parents=True, exist_ok=True)
    args.prepared_manifest.write_text(json.dumps(prepared_manifest, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    report["prepared_figure_manifest"] = str(args.prepared_manifest)

    if duplicate_sources:
        report["limitations"].append("repeated source slides are not supported in this first stable builder")
    warning_count = (
        len(report["overflow_warnings"])
        + len(report["ellipsis_warnings"])
        + len(report["image_frame_warnings"])
        + len(report["image_readability_warnings"])
        + len(report["untouched_replaceable_slots"])
        + len(report["limitations"])
    )
    report["status"] = "needs_review" if warning_count else "passed"
    report["warning_count"] = warning_count
    return report


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("plan", type=Path, help="fallback_edit_plan.json")
    parser.add_argument("output", type=Path, help="output fallback PPTX")
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--manifest", type=Path, default=DEFAULT_MANIFEST)
    parser.add_argument("--report", type=Path)
    parser.add_argument("--prepared-dir", type=Path)
    parser.add_argument("--prepared-manifest", type=Path)
    parser.add_argument("--strict", action="store_true", help="fail on expected_text mismatch")
    parser.add_argument("--fail-on-warnings", action="store_true", help="return non-zero when the build report needs review")
    args = parser.parse_args()

    if args.report is None:
        args.report = args.output.with_suffix(".build-report.json")
    if args.prepared_dir is None:
        args.prepared_dir = args.output.parent / "prepared_figures"
    if args.prepared_manifest is None:
        args.prepared_manifest = args.output.parent / "prepared_figure_manifest.json"

    try:
        report = apply_plan(args)
    except Exception as exc:
        print(f"ERROR: {exc}", file=sys.stderr)
        return 2

    args.report.parent.mkdir(parents=True, exist_ok=True)
    args.report.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(args.output)
    print(args.report)
    if args.fail_on_warnings and report["status"] != "passed":
        print(f"ERROR: build completed but report status is {report['status']}", file=sys.stderr)
        return 3
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
