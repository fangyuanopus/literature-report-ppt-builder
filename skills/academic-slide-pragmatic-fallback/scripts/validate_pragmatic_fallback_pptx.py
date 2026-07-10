#!/usr/bin/env python3
"""Structural QA for the clean pragmatic fallback route."""

from __future__ import annotations

import argparse
import json
import math
import re
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


PT_PER_CM = 28.3465
EMU_PER_CM = 360000
LINE_HEIGHT = 1.2
FONT = "Microsoft YaHei"
EXPECTED_FONT_SIZES = {
    "AGENT_NAV_LABEL_": {12},
    "AGENT_FOOTER": {8},
    "AGENT_PAGE_NUMBER": {8},
    "AGENT_TITLE": {24, 30},
    "AGENT_COVER_ENGLISH_TITLE": {20},
    "AGENT_COVER_TITLE": {30},
    "AGENT_COVER_META": {16},
    "AGENT_CLOSING": {54},
    "AGENT_CLOSING_META": {13},
    "AGENT_TAKEAWAY": {15, 16},
    "AGENT_BODY": {17, 18, 20},
    "AGENT_FIGURE_CAPTION_": {9},
    "AGENT_STEP_TEXT_": {18},
    "AGENT_STEP_ARROW_": {20},
    "AGENT_COLUMN_HEAD_": {21},
    "AGENT_SUMMARY_LABEL_": {19},
    "AGENT_SUMMARY_VALUE_": {18},
}


def visual_width(text: str) -> float:
    width = 0.0
    for char in text:
        if "\u4e00" <= char <= "\u9fff" or "\u3000" <= char <= "\u303f" or "\uff00" <= char <= "\uffef":
            width += 1.0
        elif char == " ":
            width += 0.35
        elif char.isascii():
            width += 0.5
        else:
            width += 0.8
    return width


def text_overflow(shape) -> str | None:
    """Return an overflow reason for a fixed Route C text slot, if any."""
    if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
        return None
    runs = [run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text]
    sizes = [run.font.size.pt for run in runs if run.font.size is not None]
    if not sizes:
        return None
    size = max(sizes)
    width_pt = shape.width / EMU_PER_CM * PT_PER_CM
    height_pt = shape.height / EMU_PER_CM * PT_PER_CM
    chars_per_line = max(1, math.floor(width_pt / size))
    wrap = shape.text_frame.word_wrap is not False
    max_lines = max(1, math.floor(height_pt / (size * LINE_HEIGHT))) if wrap else 1
    needed_lines = sum(max(1, math.ceil(visual_width(line) / chars_per_line)) for line in shape.text.split("\n"))
    if needed_lines > max_lines:
        return f"needs about {needed_lines} lines; slot capacity is {max_lines}"
    return None


def text_style_issues(shape) -> list[str]:
    """Enforce the fixed template's font and same-level size contract."""
    if not getattr(shape, "has_text_frame", False) or not shape.text.strip():
        return []
    matches = [(prefix, sizes) for prefix, sizes in EXPECTED_FONT_SIZES.items() if shape.name.startswith(prefix)]
    expected_sizes = max(matches, key=lambda item: len(item[0]))[1] if matches else None
    if expected_sizes is None:
        return []
    issues = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            if run.font.name != FONT:
                issues.append(f"{shape.name} uses {run.font.name or 'an inherited font'} instead of {FONT}")
            size = run.font.size.pt if run.font.size is not None else None
            if size not in expected_sizes:
                choices = "/".join(str(value) for value in sorted(expected_sizes))
                issues.append(f"{shape.name} uses {size or 'an inherited'}pt; expected {choices}pt")
    return issues


def all_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from all_shapes(shape.shapes)


def intersects(a, b) -> bool:
    return not (a.left + a.width <= b.left or b.left + b.width <= a.left or a.top + a.height <= b.top or b.top + b.height <= a.top)


def check_overlap(shapes):
    content = [s for s in shapes if any(s.name.startswith(prefix) for prefix in ("AGENT_BODY", "AGENT_FIGURE_", "AGENT_TAKEAWAY", "AGENT_COLUMN_", "AGENT_STEP_TEXT", "AGENT_SUMMARY_LABEL", "AGENT_SUMMARY_VALUE")) and "CAPTION" not in s.name]
    return [f"{a.name} overlaps {b.name}" for index, a in enumerate(content) for b in content[index + 1 :] if intersects(a, b)]


def validate(pptx: Path, plan: dict) -> dict:
    prs, expected = Presentation(pptx), plan.get("slides") or []
    issues, pages = [], []
    if len(prs.slides) != len(expected):
        issues.append(f"slide count mismatch: deck={len(prs.slides)} plan={len(expected)}")
    nav_count = len((plan.get("deck") or {}).get("navigation") or [])
    for index, slide in enumerate(prs.slides, 1):
        shapes = list(all_shapes(slide.shapes))
        names = [shape.name for shape in shapes]
        texts = [shape.text.strip() for shape in shapes if getattr(shape, "has_text_frame", False) and shape.text.strip()]
        figures = [name for name in names if name.startswith("AGENT_FIGURE_") and "CAPTION" not in name]
        expected_figures = (expected[index - 1].get("figures") or []) if index <= len(expected) else []
        expected_slide_no = (expected[index - 1].get("slide_no") if index <= len(expected) else None)
        expected_layout = (expected[index - 1].get("layout_type") if index <= len(expected) else None)
        if not expected_layout and index == 1:
            expected_layout = "cover"
        chromeless = expected_layout in ("cover", "closing")
        page_number_shapes = [shape for shape in shapes if shape.name == "AGENT_PAGE_NUMBER"]
        nav_label_count = len([name for name in names if name.startswith("AGENT_NAV_LABEL_")])
        if chromeless:
            if page_number_shapes or nav_label_count:
                issues.append(f"slide {index}: cover/closing page must not contain navigation or page number")
        else:
            if len(page_number_shapes) != 1 or not getattr(page_number_shapes[0], "has_text_frame", False):
                issues.append(f"slide {index}: missing or duplicated page number")
            elif page_number_shapes[0].text.strip() != str(expected_slide_no):
                issues.append(f"slide {index}: page number does not match plan")
            if nav_label_count != nav_count:
                issues.append(f"slide {index}: duplicated or incomplete navigation")
        if len(figures) != len(expected_figures):
            issues.append(f"slide {index}: figures deck={len(figures)} plan={len(expected_figures)}")
        unowned = [name for name in names if not name.startswith("AGENT_")]
        if unowned:
            issues.append(f"slide {index}: unowned shapes: {', '.join(unowned[:4])}")
        for shape in shapes:
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                issues.append(f"slide {index}: outside canvas: {shape.name}")
            overflow = text_overflow(shape)
            if overflow:
                issues.append(f"slide {index}: text overflow in {shape.name}: {overflow}")
            for style_issue in text_style_issues(shape):
                issues.append(f"slide {index}: style contract violation: {style_issue}")
        for text in texts:
            if re.search(r"\b(?:fig\.?\s*\d+|route\s*[abc]|image2|fallback)\b", text, re.I) and not figures:
                issues.append(f"slide {index}: unsupported visible internal/figure reference: {text[:80]}")
        for overlap in check_overlap(shapes):
            issues.append(f"slide {index}: {overlap}")
        pages.append({"slide_no": index, "shape_count": len(shapes), "figure_count": len(figures)})
    return {"$schema": "academic-pragmatic-fallback-audit/v2", "pptx": str(pptx), "status": "failed" if issues else "passed", "issues": issues, "slides": pages}


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--plan", type=Path, required=True)
    parser.add_argument("--out", type=Path)
    parser.add_argument("--fail-on-review", action="store_true")
    args = parser.parse_args()
    result = validate(args.pptx, json.loads(args.plan.read_text(encoding="utf-8")))
    data = json.dumps(result, ensure_ascii=False, indent=2) + "\n"
    if args.out:
        args.out.write_text(data, encoding="utf-8")
    else:
        print(data, end="")
    return 0 if result["status"] == "passed" else 2


if __name__ == "__main__":
    raise SystemExit(main())
